"""Phase 11 v2 — AI Vision content summary.

Runs after upload as a BackgroundTask (mirroring the Pass B face scan).
Per-category dispatch:

    image    → Florence-2 detailed caption (+ <OCR> for whiteboard /
               document / screenshot scenes) → Qwen2.5-Instruct rewriter
               composes one natural search-friendly sentence using the
               caption, named people, OCR text, and scene metadata.
    video    → ffmpeg single keyframe → image branch.
    document → server-side text extraction (pypdf / python-docx / openpyxl /
               plaintext) → BART abstractive (sumy LSA fallback).

All three paths populate `images.summary`, `summary_topic`,
`summary_points`, `summary_generated_at`, and flip `pending_summary =
false` in one update.

Each helper is best-effort: optional deps that fail to import (Florence-2
without `[ml]` extras, Qwen2.5 weights uncached, sumy without nltk data)
fall through to a regex-based deterministic path so a torch-less test
environment still produces a usable summary.
"""

from __future__ import annotations

import asyncio
import logging
import re
import subprocess
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from typing import Optional
from uuid import UUID

from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

from backend.config import settings
from backend.image import fetch_original, fetch_served
from backend.models import Face, FaceDetection, Image, Person

logger = logging.getLogger(__name__)


# --- public dataclass ------------------------------------------------------


@dataclass
class SummaryResult:
    topic: str
    summary: str
    points: list[str]
    # Optional category written to `Image.content_type` when the
    # summarizer can infer one — populated for video / document /
    # audio rows by the rule-based classifier below. Images get
    # their content_type from the vision pipeline upstream so leave
    # it None here and the existing image path's `content_type`
    # column-write isn't disturbed.
    content_type: Optional[str] = None


# --- public entry point ----------------------------------------------------


async def summarize_image_id(session: AsyncSession, image_id: UUID) -> None:
    """Load the image row, run the right strategy, write the summary back.

    Always overwrites the existing summary — switching captioning models or
    upgrading the doc extractor needs the row to be regenerable. Callers
    that don't want to redo work should check `pending_summary` themselves
    before invoking. The upload path only invokes this once per upload.
    """
    if not settings.summarize_enabled:
        return

    image = await session.get(Image, image_id)
    if image is None or image.deleted_at is not None:
        return

    # Prefer the original bytes when they're still around (highest
    # fidelity for vision models). For video rows whose original has
    # already been dropped by the transcode pipeline (per the "only
    # keep the served copy" policy), fall back to the served MP4 —
    # ffmpeg can extract keyframes from either, and the served
    # variant is the only copy that exists. Documents similarly
    # use whatever the row has.
    try:
        if image.original_blob_key is not None:
            raw_bytes, _mime = await fetch_original(image)
        elif image.served_blob_key is not None:
            raw_bytes, _mime = await fetch_served(image)
        else:
            raise RuntimeError("image row has neither original nor served bytes")
    except Exception:
        logger.exception(
            "summarize: failed to fetch source bytes for %s", image_id
        )
        await _mark_done(image_id, None)
        return

    # Look up named people while we're still in the async event loop —
    # `_dispatch` runs in a thread and can't issue DB queries against the
    # async engine. Empty list when no consent / no faces yet / no
    # named persons in this image.
    named_people: list[str] = []
    if image.category == "image":
        named_people = await _load_named_people(
            session, image.id, image.user_id
        )

    # Pre-load the existing tag labels in async land for the same
    # reason. The previous code accessed `image.tags` inside the sync
    # `_dispatch` thread; lazy-loading a relationship there blows up
    # with "greenlet_spawn has not been called" and leaves the session
    # in a PendingRollback state that makes `_mark_done` fail too.
    # Materializing the labels upfront keeps the sync path read-only.
    pre_tag_labels = await _load_image_tag_labels(session, image.id)

    # Route through the dedicated single-thread ML executor so two
    # concurrent summarize calls serialize instead of fighting for the
    # GIL with the asyncio event loop. With the default thread pool
    # the user couldn't log in while a backfill was running because
    # every request handler queued behind the model threads.
    from backend.vision.inference_pool import run_in_inference_pool
    try:
        result = await run_in_inference_pool(
            _dispatch, image, raw_bytes, named_people, pre_tag_labels
        )
    except Exception:
        logger.exception("summarize: dispatch failed for %s", image_id)
        result = None

    # Capture summary_signals BEFORE awaiting anything else — the
    # `image` object's session may be poisoned by a greenlet-less
    # lazy load that happened inside _dispatch; reading attributes
    # later goes through state-load paths that re-raise. We grab the
    # in-memory dict (set by `_summarize_image`) now, while we're
    # still in the same thread context, and pass it positionally
    # to `_mark_done` which uses a fresh session for the UPDATE.
    signals = None
    try:
        signals = image.__dict__.get("summary_signals")
    except Exception:
        pass

    # Capture the owner id now (same safe `__dict__` read used for
    # doc-chunk persistence below) so `_mark_done` can scope the
    # adjective-tag write to this user without re-touching the
    # possibly-poisoned ORM state.
    owner_id = None
    try:
        owner_id = image.__dict__.get("user_id")
    except Exception:
        owner_id = None

    try:
        await _mark_done(image_id, result, signals, owner_id)
    except Exception:
        logger.exception("summarize: _mark_done failed for %s", image_id)

    # Sprint I D2 — persist document chunks for jump-to-section search.
    # Only relevant for document rows; image / video paths don't stash
    # `doc_chunks`. Runs OUTSIDE the inference pool because it's just
    # database I/O at this point — the embedding work happened inside
    # the sync dispatch already.
    doc_chunks = None
    try:
        doc_chunks = image.__dict__.get("doc_chunks")
    except Exception:
        doc_chunks = None
    if doc_chunks:
        try:
            await _persist_doc_chunks(image_id, image.user_id, doc_chunks)
        except Exception:
            logger.exception(
                "summarize: doc chunk persistence failed for %s", image_id
            )


async def _persist_doc_chunks(
    image_id, user_id, chunks: list[dict],
) -> None:
    """Idempotent UPSERT of per-chunk doc embeddings.

    Uses ON CONFLICT (image_id, chunk_index) DO UPDATE so re-
    summarization (admin backfill / model swap) overwrites in place
    instead of duplicating. Embedding can be None — the search path
    will fall back to FTS for rows with no vector.

    Runs in a fresh session — same pattern as `_mark_done` — to side-
    step any greenlet poisoning from the sync dispatch thread.
    """
    if not chunks:
        return
    from sqlalchemy.dialects.postgresql import insert as pg_insert
    from backend.db import SessionLocal
    from backend.models import DocumentChunk

    rows = [
        {
            "image_id": image_id,
            "user_id": user_id,
            "chunk_index": c["chunk_index"],
            "text": c["text"],
            "embedding": c.get("embedding"),
        }
        for c in chunks
    ]
    async with SessionLocal() as session:
        stmt = pg_insert(DocumentChunk).values(rows)
        stmt = stmt.on_conflict_do_update(
            index_elements=["image_id", "chunk_index"],
            set_={
                "text": stmt.excluded.text,
                "embedding": stmt.excluded.embedding,
            },
        )
        await session.execute(stmt)
        await session.commit()


async def _load_image_tag_labels(
    session: AsyncSession, image_id
) -> list[str]:
    """Return the labels on this image's existing tags.

    Run from the async side so the sync `_dispatch` thread doesn't have
    to lazy-load `image.tags` itself (which fails with a greenlet
    error). Caller passes the resulting list to `_summarize_image` as
    `pre_tag_labels`.
    """
    from backend.models import ImageTag, Tag
    rows = (
        await session.execute(
            select(Tag.label)
            .join(ImageTag, ImageTag.tag_id == Tag.id)
            .where(ImageTag.image_id == image_id)
        )
    ).scalars().all()
    return [r for r in rows if r]


async def _load_named_people(
    session: AsyncSession, image_id, user_id
) -> list[str]:
    """Distinct, ordered names of identified people in this image.

    Pass B (faces) populates `face_detections.face_id → faces.person_id →
    persons.display_name`. Anonymous detections (face exists but no person
    named yet) are filtered out — generic captions are better than ones
    that say "and another person" with no anchor.

    §D1 (Sprint I) — when a person record is labelled "Me" / "I" /
    "myself" (the user's own face cluster, which face-recognition UX
    encourages people to tag themselves as), substitute the user's
    actual display_name from the users table. This is what turns
    summaries like "Me holding a coffee" into "Jason holding a
    coffee" — the second one indexes correctly and reads as the
    user expects when they search for their own name.
    """
    # Owner's display name. Loaded once up-front so the "Me" → name
    # swap below can run inline. None / empty leaves the "Me" label
    # intact (existing behavior; first-person pronoun polish handles
    # downstream pronouns).
    from backend.models import User

    owner_display = (
        await session.execute(
            select(User.display_name).where(User.id == user_id)
        )
    ).scalar_one_or_none()
    owner_display = (owner_display or "").strip() or None

    # Min(bbox_x) lets us order persons left-to-right even with DISTINCT —
    # plain ORDER BY bbox_x violates the SELECT-list rule on Postgres.
    rows = (
        await session.execute(
            select(Person.display_name, func.min(FaceDetection.bbox_x).label("x"))
            .join(Face, Face.person_id == Person.id)
            .join(FaceDetection, FaceDetection.face_id == Face.id)
            .where(
                FaceDetection.image_id == image_id,
                FaceDetection.user_id == user_id,
                Person.display_name.is_not(None),
                Person.display_name != "",
            )
            .group_by(Person.display_name)
            .order_by("x")
        )
    ).all()

    first_person_aliases = {"me", "i", "myself"}
    out: list[str] = []
    for r in rows:
        name = (r[0] or "").strip()
        if not name:
            continue
        if owner_display and name.lower() in first_person_aliases:
            name = owner_display
        out.append(name)
    return out


async def _mark_done(
    image_id,
    result: Optional[SummaryResult],
    signals: Optional[dict] = None,
    user_id=None,
) -> None:
    """Persist the summary row after a worker run.

    Takes only the id + flat data (no ORM-tracked Image instance).
    Accessing attributes on an Image object that's bound to a
    poisoned session re-triggers the same PendingRollbackError on a
    fresh session, because the attribute access goes through the
    InstanceState's load_expired/refresh path. Writing via a plain
    `UPDATE ... WHERE id = :id` on a fresh session sidesteps that
    entirely.

    When `result is None` the dispatch produced nothing usable — the LLM
    crashed, the model wasn't loadable, the file was corrupt, etc. In
    that case we deliberately keep `pending_summary=True` so the row
    stays visible to the regular backfill pass (which targets
    `pending_summary=true OR summary IS NULL`). Marking it complete
    would silently drop the row from the queue and leave the user with
    no summary forever.
    """
    from sqlalchemy import update as sa_update

    from backend.db import SessionLocal

    values: dict = {}
    if result is not None:
        values["summary"] = result.summary
        values["summary_topic"] = result.topic
        values["summary_points"] = result.points
        values["pending_summary"] = False
        values["summary_generated_at"] = datetime.now(timezone.utc)
        if signals:
            values["summary_signals"] = signals
        if result.content_type:
            values["content_type"] = result.content_type
        # Encode the summary in CLIP text space so search can score
        # against text-to-text semantic distance (not just image-to-text
        # visual cosine). Topic is prepended because it often carries
        # the noun phrase that anchors the semantics ("Cat Sleeping
        # On Laptop Keyboard" vs. the longer prose summary).
        try:
            embedding = await asyncio.to_thread(
                _encode_summary_for_search, result.summary, result.topic,
            )
        except Exception:
            logger.exception("summary embed: to_thread failed")
            embedding = None
        if embedding is not None:
            values["summary_clip_embedding"] = embedding
    else:
        values["summary_generated_at"] = datetime.now(timezone.utc)
        # Keep pending_summary alone so the row gets retried.

    # SYNC WRITE on purpose. The previous `async with SessionLocal()`
    # path opens a connection from the process-wide async pool, which
    # is bound to whatever event loop FIRST created it. When the ML
    # worker's heartbeat task crashes (asyncpg's "Task got Future
    # attached to a different loop" error pattern), the pool's
    # connection futures get poisoned. The next `_mark_done` call
    # from a different task inherits the bad state and the row write
    # silently fails — summarize ran, summary was generated, but the
    # row stays at `pending_summary=true` forever.
    #
    # psycopg2 sync connection sidesteps all of that: own its own
    # connection, no event-loop binding, no shared pool. The write
    # is one round-trip so the blocking call is cheap (~5 ms).
    # `asyncio.to_thread` keeps the event loop responsive while we
    # do it, same pattern as the embedding call above.
    await asyncio.to_thread(_mark_done_sync, image_id, values)

    # --- Adjective tags from the just-written summary -----------------
    # The row's tags should be the descriptive adjectives in its summary
    # (user request) rather than the upload-time CLIP concept labels.
    # Best-effort and fully isolated: any failure here logs and leaves
    # the prior tags intact — it must NEVER break the summarize path.
    # Only runs when we produced a real summary and know the owner.
    if result is not None and user_id is not None and result.summary:
        try:
            adjectives = await asyncio.to_thread(
                _extract_adjective_tags, result.summary
            )
            if adjectives:
                await asyncio.to_thread(
                    _write_adjective_tags_sync, image_id, user_id, adjectives,
                )
                logger.info(
                    "adjective tags for %s: %s", image_id, ", ".join(adjectives)
                )
            else:
                logger.info(
                    "adjective tags for %s: none extracted; keeping prior tags",
                    image_id,
                )
        except Exception:
            logger.exception("adjective tags: write failed for %s", image_id)


def _mark_done_sync(image_id, values: dict) -> None:
    """Synchronous row update — runs in a thread to avoid the
    asyncpg-pool loop-binding issue described in `_mark_done`.
    psycopg2 dependency is already in the base deps for Alembic;
    no new package required."""
    import json
    import psycopg2

    sync_url = settings.database_url_sync.replace(
        "postgresql+psycopg2://", "postgresql://"
    )
    cols: list[str] = []
    params: list = []
    for k, v in values.items():
        if k == "summary_points":
            cols.append(f"{k} = %s::jsonb")
            params.append(json.dumps(v) if v is not None else None)
        elif k == "summary_signals":
            cols.append(f"{k} = %s::jsonb")
            params.append(json.dumps(v) if v is not None else None)
        elif k == "summary_clip_embedding":
            cols.append(f"{k} = %s::vector")
            params.append(v)
        else:
            cols.append(f"{k} = %s")
            params.append(v)
    params.append(str(image_id))
    sql = f"UPDATE images SET {', '.join(cols)} WHERE id = %s"
    conn = psycopg2.connect(sync_url)
    try:
        with conn.cursor() as cur:
            cur.execute(sql, params)
        conn.commit()
    finally:
        conn.close()


# ---------- Content-type classification (rule-based) -----------------
#
# Maps the AI summary + filename + extension to one of a small fixed
# taxonomy. The values land in `Image.content_type` so the gallery
# can facet by them and search can match queries like "find my
# receipts" or "show me tutorials" against this column directly.
#
# Rule-based on purpose: a separate Qwen round-trip per upload would
# double the summarize budget for a benefit that mostly disappears
# behind keyword overlap. Easy to evolve into LLM-prompted
# classification later if precision matters more than latency.

_VIDEO_CATEGORIES: dict[str, list[str]] = {
    # category → marker phrases (lowercase, substring match against summary)
    "tutorial":    ["tutorial", "how to", "step-by-step", "guide", "walkthrough", "lesson"],
    "screencast":  ["screen recording", "desktop", "screencast", "tutorial recording"],
    "recipe":      ["recipe", "cooking", "kitchen", "baking", "ingredients", "stove", "oven"],
    "sports":      ["game", "match", "ball", "field", "court", "stadium", "running", "swimming"],
    "music":       ["concert", "band", "musician", "guitar", "drums", "singing", "performance"],
    "gaming":      ["gameplay", "gaming", "controller", "console", "minecraft", "fortnite"],
    "vlog":        ["vlog", "talking", "selfie", "blogger", "speaking to camera"],
    "family":      ["family", "child", "kids", "baby", "birthday", "wedding", "anniversary"],
    "travel":      ["travel", "vacation", "trip", "mountain", "beach", "city skyline"],
    "animation":   ["animation", "animated", "cartoon", "3d render"],
    "presentation":["slide", "presentation", "powerpoint", "lecture", "speaker"],
}

_DOC_CATEGORIES: dict[str, list[str]] = {
    "code":         [],  # by extension (see DOC_EXT_CATEGORIES)
    "spreadsheet":  [],
    "presentation":[],
    "contract":     ["contract", "agreement", "lease", "terms and conditions", "party of the first", "hereby"],
    "receipt":      ["receipt", "invoice", "total due", "subtotal", "transaction", "purchase"],
    "recipe":       ["recipe", "ingredients", "preheat", "tablespoon", "cooking instruction"],
    "manual":       ["manual", "user guide", "installation", "specifications", "warranty"],
    "report":       ["report", "executive summary", "findings", "methodology", "abstract"],
    "letter":       ["dear", "sincerely", "regards", "letter", "memo"],
    "form":         ["please fill", "checkbox", "questionnaire", "form", "applicant"],
    "notes":        ["notes", "todo", "to-do", "journal", "diary"],
    "legal":        ["court", "plaintiff", "defendant", "statute", "ordinance"],
    "research":     ["abstract", "hypothesis", "references", "doi:", "citation", "et al."],
}

_DOC_EXT_CATEGORIES: dict[str, str] = {
    "py": "code", "js": "code", "ts": "code", "tsx": "code", "jsx": "code",
    "rs": "code", "go": "code", "rb": "code", "java": "code", "kt": "code",
    "c": "code", "cpp": "code", "h": "code", "hpp": "code", "cs": "code",
    "swift": "code", "m": "code", "mm": "code", "sh": "code", "ps1": "code",
    "sql": "code", "r": "code", "lua": "code", "scala": "code",
    "xlsx": "spreadsheet", "xls": "spreadsheet", "csv": "spreadsheet",
    "tsv": "spreadsheet", "ods": "spreadsheet",
    "pptx": "presentation", "ppt": "presentation", "odp": "presentation",
    "key": "presentation",
}


def _classify_content(
    summary: str | None,
    filename: str | None,
    kind: str,  # "video" | "document" | "audio"
) -> Optional[str]:
    """Return a content_type label for the row, or None when no
    rule fires. `kind` selects the taxonomy. Filename extension is
    consulted first when it's a strong signal (code / spreadsheet /
    presentation) — extension trumps body content in those cases
    because a .py file is code regardless of what its contents
    parse as in plain English.
    """
    if kind == "document" and filename:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in _DOC_EXT_CATEGORIES:
            return _DOC_EXT_CATEGORIES[ext]
    table = _VIDEO_CATEGORIES if kind == "video" else _DOC_CATEGORIES
    text = (summary or "").lower()
    if not text:
        return None
    # Score each category by how many marker phrases it hits;
    # winner takes the top, tie → first in insertion order.
    best_cat = None
    best_hits = 0
    for cat, markers in table.items():
        if not markers:
            continue
        hits = sum(1 for m in markers if m in text)
        if hits > best_hits:
            best_cat = cat
            best_hits = hits
    return best_cat


def _encode_summary_for_search(
    summary: str | None, topic: str | None,
) -> list[float] | None:
    """Run the summary + topic through the CLIP text encoder and
    return a list[float] suitable for pgvector. Returns None on any
    failure (ml extras not installed, encoder OOM, empty text) so
    `_mark_done` can carry on with a NULL embedding — the search
    path treats NULL as "no text-side semantic signal" and falls
    back to image-cosine + FTS, which is the prior behavior.
    """
    parts = [s for s in (topic, summary) if s and s.strip()]
    if not parts:
        return None
    text = " — ".join(parts)
    # Truncate to a safe token budget so the encoder doesn't reject
    # very long summaries. CLIP ViT-L-14 tokenizer caps at 77 tokens;
    # ~250 characters is well under that for English text.
    if len(text) > 250:
        text = text[:250]
    try:
        from backend.vision.runtime import encode_text_cached
        vec = encode_text_cached(text)
    except Exception:
        logger.exception("summary embed failed for text=%r", text[:60])
        return None
    return [float(x) for x in vec]


# ---------- Adjective tags from the generated summary ----------------
#
# The user's file tags should be the DESCRIPTIVE ADJECTIVES that appear
# in the file's AI summary ("serene", "vibrant", "dark", "minimalist"),
# not the CLIP object/concept labels ("dog", "kitchen", "selfie") the
# upload pass writes. After a summary is generated we POS-tag it,
# keep the adjective tokens (Penn-Treebank JJ / JJR / JJS), and write
# those as the row's tags — replacing the upload-time CLIP tags.
#
# Extraction uses NLTK (already a hard dep of the `[ml]` extras — it
# backs `sumy` document summarization and `synonyms.py` WordNet). The
# `averaged_perceptron_tagger_eng` + `punkt_tab` data download lazily
# on first use (and are baked into the Dockerfile for reliability); a
# failed load is cached so we don't re-pay the cost, and the caller
# falls back to keeping the existing tags rather than crashing.

# Adjectives that are grammatically JJ but carry no descriptive value as
# a tag — determiners / quantifiers / deictics the tagger labels JJ, plus
# a few nouns Penn-Treebank routinely mis-tags as JJ when they sit in an
# attributive slot ("video game" → "video"/JJ, "self-portrait" → JJ).
_TAG_ADJ_STOPWORDS: frozenset[str] = frozenset({
    # determiner / quantifier / deictic adjectives
    "this", "that", "these", "those", "other", "another", "such", "same",
    "own", "more", "most", "many", "much", "few", "fewer", "little", "less",
    "least", "several", "various", "certain", "whole", "entire", "overall",
    "only", "single", "double", "multiple", "first", "second", "third",
    "last", "next", "previous", "former", "latter", "new", "old", "good",
    "great", "real", "sure", "able", "due", "likely", "available",
    # nouns the POS tagger commonly mis-labels JJ in attributive position
    "video", "image", "photo", "picture", "self-portrait", "self",
    "close-up", "screenshot", "screen", "front", "back", "top", "left",
    "right", "side", "today", "tomorrow",
})

# Sentinel for the lazy NLTK-data load: None = not tried, True = ready,
# False = tried and failed (fall back to the regex path).
_NLTK_POS_READY: Optional[bool] = None


def _ensure_nltk_pos() -> bool:
    """Make sure NLTK's POS tagger + tokenizer data are loadable.

    Mirrors the lazy-load pattern in `backend/synonyms.py`: try to use
    the resource, download it once if missing, cache the outcome. The
    Dockerfile pre-downloads these so the first request doesn't pay the
    network cost — this runtime path is the belt-and-suspenders fallback
    for dev shells / fresh installs. Returns False (cached) when NLTK
    isn't installed or the data can't be fetched, so callers degrade to
    the regex extractor instead of raising.
    """
    global _NLTK_POS_READY
    if _NLTK_POS_READY is not None:
        return _NLTK_POS_READY
    try:
        import nltk  # type: ignore

        # punkt_tab (tokenizer) + averaged_perceptron_tagger_eng (POS).
        # NLTK >= 3.9 renamed both resources with the `_tab` / `_eng`
        # suffixes; we target those names explicitly. `find` raises
        # LookupError when the data is absent — that's our download cue.
        for finder, pkg in (
            ("tokenizers/punkt_tab", "punkt_tab"),
            ("taggers/averaged_perceptron_tagger_eng",
             "averaged_perceptron_tagger_eng"),
        ):
            try:
                nltk.data.find(finder)
            except LookupError:
                nltk.download(pkg, quiet=True)
        # Smoke-test the full path so a half-installed corpus fails here,
        # not mid-summary.
        from nltk import pos_tag, word_tokenize
        pos_tag(word_tokenize("a quick brown test"))
        _NLTK_POS_READY = True
    except Exception as e:
        logger.info("adjective tags: NLTK POS unavailable (%s); regex fallback", e)
        _NLTK_POS_READY = False
    return _NLTK_POS_READY


def _clean_adj_token(word: str) -> Optional[str]:
    """Lower-case + validate a candidate adjective token.

    Keeps single hyphens inside the word (so "shoulder-length",
    "ice-themed", "black-and-white" survive) but drops anything with
    digits, punctuation, or that is too short. Returns None when the
    token isn't a usable tag.
    """
    w = word.lower().strip().strip("-")
    # letters + internal single hyphens only, length >= 3
    if not re.fullmatch(r"[a-z]+(?:-[a-z]+)*", w):
        return None
    if len(w) < 3:
        return None
    return w


def _extract_adjectives_regex(summary: str, cap: int) -> list[str]:
    """Crude adjective heuristic for when NLTK data isn't available.

    Looks for words ending in common adjective-forming suffixes
    (-ful, -ous, -ive, -ish, -less, -ic, -al, -y, …). Lower precision
    than POS tagging but keeps the feature alive on a torch-less /
    data-less box instead of returning nothing.
    """
    suffix = re.compile(
        r".+(?:ful|ous|ive|ish|less|able|ible|ic|ical|al|ial|"
        r"ant|ent|ary|y|y|ed)$"
    )
    out: list[str] = []
    seen: set[str] = set()
    for raw in re.findall(r"[A-Za-z][A-Za-z-]*[A-Za-z]", summary):
        w = _clean_adj_token(raw)
        if not w or w in seen or w in _TAG_ADJ_STOPWORDS:
            continue
        if not suffix.match(w):
            continue
        seen.add(w)
        out.append(w)
        if len(out) >= cap:
            break
    return out


def _extract_adjective_tags(summary: str | None, cap: int = 8) -> list[str]:
    """Return up to `cap` descriptive adjectives from `summary`.

    Deduped, lower-cased, ordered by first appearance in the text
    (salience proxy — a summary leads with its most defining traits).
    POS-tags via NLTK and keeps JJ / JJR / JJS tokens; falls back to a
    suffix-based regex heuristic when NLTK data can't load. Returns an
    empty list for empty / very short summaries so the caller can keep
    the previous tags instead of wiping them.
    """
    if not summary or not summary.strip():
        return []
    text = summary.strip()

    if _ensure_nltk_pos():
        try:
            from nltk import pos_tag, word_tokenize

            out: list[str] = []
            seen: set[str] = set()
            for word, tag in pos_tag(word_tokenize(text)):
                if tag not in ("JJ", "JJR", "JJS"):
                    continue
                w = _clean_adj_token(word)
                if not w or w in seen or w in _TAG_ADJ_STOPWORDS:
                    continue
                seen.add(w)
                out.append(w)
                if len(out) >= cap:
                    break
            return out
        except Exception:
            logger.exception("adjective tags: POS extraction failed; regex fallback")

    return _extract_adjectives_regex(text, cap)


def _write_adjective_tags_sync(image_id, user_id, adjectives: list[str]) -> None:
    """Replace this image's auto-derived tags with `adjectives`.

    Synchronous psycopg2 write, same rationale as `_mark_done_sync`:
    owns its own connection (no asyncpg event-loop binding) and runs
    inside `asyncio.to_thread`. Connects as the `neuthek` superuser,
    which bypasses the FORCE-RLS policies on `tags` / `image_tags`.

    Semantics:
      * Detach every CURRENT image_tags row for this image whose tag
        was machine-generated (`source IN ('clip','auto')`). User-
        applied tags (`source='user'`) are left untouched — the user
        curated those by hand.
      * Upsert a Tag row per adjective (per-user, case-folded unique,
        `source='auto'`) and link it to the image. Dedupe is on
        `lower(label)` to match the `tags_user_label_idx` functional
        unique index.

    No-op when `adjectives` is empty (caller guarantees this only fires
    when extraction produced something), so a row with an un-taggable
    summary keeps whatever tags it already had.
    """
    if not adjectives:
        return
    import psycopg2

    sync_url = settings.database_url_sync.replace(
        "postgresql+psycopg2://", "postgresql://"
    )
    conn = psycopg2.connect(sync_url)
    try:
        with conn.cursor() as cur:
            # 1. Drop existing machine tags from THIS image only. Other
            #    images sharing the same Tag row keep their links; the
            #    Tag row itself is left in place (cheap, and re-used on
            #    the next image that needs it).
            cur.execute(
                """
                DELETE FROM image_tags it
                USING tags t
                WHERE it.image_id = %s
                  AND it.tag_id = t.id
                  AND t.source IN ('clip', 'auto')
                """,
                (str(image_id),),
            )

            # 2. Upsert each adjective Tag (per-user) and link it.
            for label in adjectives:
                # Resolve / create the tag row, case-folded per user.
                cur.execute(
                    "SELECT id FROM tags "
                    "WHERE user_id = %s AND lower(label) = lower(%s) "
                    "LIMIT 1",
                    (str(user_id), label),
                )
                row = cur.fetchone()
                if row:
                    tag_id = row[0]
                else:
                    cur.execute(
                        "INSERT INTO tags (user_id, label, source) "
                        "VALUES (%s, %s, 'auto') RETURNING id",
                        (str(user_id), label),
                    )
                    tag_id = cur.fetchone()[0]

                # Link to the image (idempotent — PK is (image_id, tag_id)).
                cur.execute(
                    "INSERT INTO image_tags (image_id, tag_id, user_id) "
                    "VALUES (%s, %s, %s) "
                    "ON CONFLICT (image_id, tag_id) DO NOTHING",
                    (str(image_id), tag_id, str(user_id)),
                )
        conn.commit()
    finally:
        conn.close()


def _dispatch(
    image: Image,
    raw_bytes: bytes,
    named_people: list[str],
    pre_tag_labels: list[str] | None = None,
) -> Optional[SummaryResult]:
    if image.category == "image":
        return _summarize_image(image, raw_bytes, named_people, pre_tag_labels or [])
    if image.category == "video":
        return _summarize_video(image, raw_bytes)
    if image.category == "audio":
        return _summarize_audio(image, raw_bytes)
    if image.category == "document":
        return _summarize_document(image, raw_bytes)
    # Catch-all for "other" — archives (.zip, .tar.gz, etc.) and any
    # mime type we don't have a dedicated summarizer for. Without this,
    # the row's `summary` stays NULL forever and the progress banner
    # is stuck at N-1 of N. Archives now describe their inner listing
    # (#184); anything unreadable falls back to a filename-derived stub.
    return _summarize_other(image, raw_bytes)


def _summarize_other(image: Image, raw_bytes: bytes | None = None) -> SummaryResult:
    """Summarize an "other" blob — archives (.zip/.tar/.tar.gz) and any
    mime without a dedicated summarizer.

    #184 — archives now describe their CONTENTS, not just the filename.
    We read the inner listing (zip via stdlib zipfile, tar/tgz via
    tarfile) and infer what the archive HOLDS — "Zip archive of a React
    project", "Zip archive of 40 photos", "Tar archive of Python source
    (12 files)" — so semantic search can surface it by what's inside.
    Falls back to the filename-derived stub when we can't read the
    listing (encrypted zip, unknown format, bytes unavailable).
    """
    fname = image.original_filename or "file"
    lower = fname.lower()
    ext = (fname.rsplit(".", 1)[-1].lower() if "." in fname else "") or ""
    # Stem for the "Named '<x>'" suffix — strip a compound archive
    # extension (.tar.gz / .tar.bz2) fully so we don't show "project.tar".
    stem_src = fname
    for compound in (".tar.gz", ".tar.bz2", ".tar.xz", ".tgz", ".tbz2"):
        if lower.endswith(compound):
            stem_src = fname[: -len(compound)]
            break
    else:
        stem_src = fname.rsplit(".", 1)[0] if "." in fname else fname
    stem = stem_src.replace("_", " ").replace("-", " ").strip()
    # Compound tar extensions read as "Tar archive" regardless of the
    # final compression suffix (.tar.gz is a tarball, not a lone gzip).
    if lower.endswith((".tar.gz", ".tar.bz2", ".tar.xz", ".tgz", ".tbz2")):
        ext_label = "Tar archive"
    else:
        ext_label = {
            "zip": "Zip archive",
            "tar": "Tar archive",
            "gz":  "Gzip archive",
            "bz2": "Bzip2 archive",
            "rar": "RAR archive",
            "7z":  "7-Zip archive",
            "epub": "EPUB e-book",
        }.get(ext, ext.upper() + " file" if ext else "File")

    names: list[str] = []
    if raw_bytes:
        try:
            if lower.endswith(".zip") or lower.endswith(".epub"):
                names = _list_zip_members(raw_bytes)
            elif (lower.endswith(".tar") or lower.endswith(".tar.gz")
                  or lower.endswith(".tgz") or lower.endswith(".tar.bz2")
                  or lower.endswith(".gz") or lower.endswith(".bz2")):
                names = _list_tar_members(raw_bytes)
        except Exception:
            logger.exception("archive listing failed for %s", fname)
            names = []

    if names:
        descr, signals = _describe_archive_contents(names)
        summary = f"{ext_label} containing {descr}."
        if stem:
            summary += f" Named '{stem}'."
        topic = (f"{ext_label}: {descr}")[:90]
        points: list[str] = [f"{len(names)} entries"]
        # Surface a few representative inner names for the preview panel
        # + so FTS indexes them.
        sample = [n for n in names if not n.endswith("/")][:6]
        if sample:
            points.append("Includes: " + ", ".join(
                n.rsplit("/", 1)[-1] for n in sample
            ))
        # Stash inner-name + project signals so search's haystack
        # (signals.concepts) indexes the archive's contents.
        image.__dict__["summary_signals"] = {
            "kind": "archive",
            "entry_count": len(names),
            "concepts": signals[:20],
        }
        return SummaryResult(
            topic=topic or ext_label, summary=summary, points=points[:5],
            content_type="archive",
        )

    # No listing — filename-only stub (prior behavior, slightly richer).
    topic = stem[:80] if stem else ext_label
    summary = f"{ext_label} — {stem}." if stem else f"{ext_label}."
    return SummaryResult(
        topic=topic or "File", summary=summary, points=[],
        content_type="archive" if ext in {
            "zip", "tar", "gz", "tgz", "bz2", "rar", "7z", "epub"
        } else None,
    )


def _list_zip_members(raw: bytes, cap: int = 2000) -> list[str]:
    """Inner file names of a zip (also used for .epub). Returns [] on any
    failure (encrypted, truncated, not a zip). Caps the list so a zip
    with 100k entries doesn't balloon memory."""
    import zipfile
    out: list[str] = []
    try:
        with zipfile.ZipFile(BytesIO(raw)) as zf:
            for info in zf.infolist():
                out.append(info.filename)
                if len(out) >= cap:
                    break
    except Exception:
        return []
    return out


def _list_tar_members(raw: bytes, cap: int = 2000) -> list[str]:
    """Inner file names of a tar / tar.gz / tar.bz2. Returns [] on any
    failure. `tarfile` auto-detects the compression from the stream."""
    import tarfile
    out: list[str] = []
    try:
        with tarfile.open(fileobj=BytesIO(raw), mode="r:*") as tf:
            for member in tf:
                out.append(member.name)
                if len(out) >= cap:
                    break
    except Exception:
        return []
    return out


# Inner-file extension → coarse content bucket, for describing what an
# archive holds ("photos", "source code", "documents", …).
_ARCHIVE_BUCKETS: tuple[tuple[str, frozenset[str]], ...] = (
    ("photos", frozenset({
        "jpg", "jpeg", "png", "gif", "heic", "heif", "webp", "bmp",
        "tiff", "raw", "cr2", "nef", "dng",
    })),
    ("videos", frozenset({"mp4", "mov", "avi", "mkv", "webm", "m4v"})),
    ("audio files", frozenset({"mp3", "m4a", "wav", "flac", "ogg", "aac"})),
    ("documents", frozenset({
        "pdf", "doc", "docx", "xls", "xlsx", "ppt", "pptx", "odt",
        "txt", "rtf", "csv", "md",
    })),
    ("source code", frozenset({
        "py", "js", "ts", "tsx", "jsx", "go", "rs", "java", "kt", "rb",
        "php", "c", "h", "cpp", "cs", "swift", "scala", "sh", "sql",
    })),
    ("images / assets", frozenset({"svg", "ico", "ttf", "woff", "woff2"})),
)

# Signature files that identify a project type from the inner listing.
_ARCHIVE_PROJECT_MARKERS: tuple[tuple[str, str], ...] = (
    ("package.json", "a Node.js / JavaScript project"),
    ("tsconfig.json", "a TypeScript project"),
    ("requirements.txt", "a Python project"),
    ("pyproject.toml", "a Python project"),
    ("setup.py", "a Python project"),
    ("cargo.toml", "a Rust project"),
    ("go.mod", "a Go project"),
    ("pom.xml", "a Java / Maven project"),
    ("build.gradle", "a Gradle project"),
    ("gemfile", "a Ruby project"),
    ("composer.json", "a PHP project"),
    ("dockerfile", "a Dockerized project"),
    ("index.html", "a website / web project"),
)


def _describe_archive_contents(names: list[str]) -> tuple[str, list[str]]:
    """Return (human description, signal-tokens) for an archive listing.

    Recognizes a project type from marker files first ("a React
    project"), otherwise describes by the dominant inner content bucket
    ("40 photos", "Python source and documents"). The signal-token list
    feeds search's concept haystack so the archive is findable by what it
    contains.
    """
    files = [n for n in names if not n.endswith("/")]
    base_lower = [n.rsplit("/", 1)[-1].lower() for n in files]

    # Project detection — react gets special-cased off package.json + a
    # jsx/tsx presence; otherwise the first matching marker wins.
    has_react = any(
        b in ("package.json",) for b in base_lower
    ) and any(
        n.lower().endswith((".jsx", ".tsx")) for n in files
    )
    project: Optional[str] = None
    if has_react:
        project = "a React project"
    else:
        marker_set = set(base_lower)
        for marker, label in _ARCHIVE_PROJECT_MARKERS:
            if marker in marker_set:
                project = label
                break

    # Content buckets by extension frequency.
    counts: dict[str, int] = {}
    for b in base_lower:
        ext = b.rsplit(".", 1)[-1] if "." in b else ""
        for bucket, exts in _ARCHIVE_BUCKETS:
            if ext in exts:
                counts[bucket] = counts.get(bucket, 0) + 1
                break
    ranked = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)

    n_files = len(files)
    signals: list[str] = []
    if project:
        signals.append(project.replace("a ", "").replace("an ", ""))

    # Compose the description.
    if project and ranked:
        top_bucket = ranked[0][0]
        descr = project
        # Add a photo/media count when that's clearly the bulk.
        if top_bucket in ("photos", "videos", "audio files") and ranked[0][1] >= 3:
            descr = f"{project} with {ranked[0][1]} {top_bucket}"
    elif project:
        descr = project
    elif ranked:
        top_bucket, top_n = ranked[0]
        # Pure single-bucket archive → "40 photos".
        if len(ranked) == 1 or ranked[0][1] >= 3 * (ranked[1][1] if len(ranked) > 1 else 1):
            descr = f"{top_n} {top_bucket}" if top_bucket in (
                "photos", "videos", "audio files"
            ) else top_bucket
        else:
            # Mixed → name the top two buckets.
            descr = " and ".join(b for b, _ in ranked[:2])
        signals.extend(b for b, _ in ranked[:3])
    else:
        descr = f"{n_files} files" if n_files else "files"

    # Always add the top few inner file extensions as signal tokens so
    # search can hit "zip with .py files".
    ext_tokens = []
    for b in base_lower:
        if "." in b:
            ext_tokens.append(b.rsplit(".", 1)[-1])
    for t in ext_tokens[:10]:
        if t and t not in signals:
            signals.append(t)

    return descr, signals


# --- image -----------------------------------------------------------------


# Scenes / content types where reading visible text materially improves
# the summary. We don't OCR every cat photo — Florence-2 OCR is ~1-2 s on
# GPU, would dominate latency on bulk re-summarize.
_OCR_SCENES = frozenset({
    "classroom", "lecture_room", "conference_room", "computer_room",
    "library", "office", "bookstore", "whiteboard",
})

_OCR_CONTENT_TYPES = frozenset({"screenshot", "document"})


def _summarize_image(
    image: Image,
    raw_bytes: bytes,
    named_people: list[str],
    pre_tag_labels: list[str] | None = None,
) -> SummaryResult:
    """Florence-2 detailed caption + scene-gated OCR + LLM rewrite.

    Flow:
      1. Florence-2 <MORE_DETAILED_CAPTION>  — dense scene description.
      2. Florence-2 <OCR>                    — only for whiteboard /
                                               classroom / screenshot etc.
      3. Qwen2.5-Instruct rewrite            — one natural sentence
                                               combining all signals.
      4. Regex fallback (clean+splice)       — when the LLM is unavailable.

    The fallback chain means the function still returns a usable summary
    when only some of the models are loadable, instead of crashing or
    emitting a generic placeholder.
    """
    raw_caption = (_caption_image(raw_bytes) or "").strip()

    needs_ocr = (
        (image.content_type or "") in _OCR_CONTENT_TYPES
        or (image.scene_label or "") in _OCR_SCENES
    )
    ocr_text = _ocr_image(raw_bytes) if needs_ocr else None

    # C2a — multi-model image pipeline. Each stage is best-effort and
    # returns None on failure so a missing model degrades the summary
    # instead of breaking it.
    regions = _florence_regions(raw_bytes)
    objects = _florence_objects(raw_bytes)
    # CLIP concept tags + heavy VLM are wired in C2b / C2e.
    concepts: Optional[list[str]] = None
    vlm_description: Optional[str] = None
    try:
        from backend.vision.concepts import top_concepts
        concepts = top_concepts(raw_bytes)
    except Exception:
        concepts = None
    try:
        if getattr(settings, "heavy_vlm_enabled", False):
            from backend.vision.runtime import get_internvl2  # noqa: F401
            vlm_description = _vlm_describe(raw_bytes)
    except Exception:
        vlm_description = None

    # Existing image_tags labels — pre-loaded in async land by
    # `summarize_image_id` so the sync thread doesn't trip on lazy
    # relationship loading.
    pre_tags: list[str] = list(pre_tag_labels or [])

    summary = _llm_rewrite_summary(
        caption=raw_caption,
        names=named_people,
        ocr_text=ocr_text,
        scene=image.scene_label,
        setting=image.indoor_outdoor,
        content_type=image.content_type,
        tags=pre_tags or None,
        regions=regions,
        objects=objects,
        concepts=concepts,
        vlm_description=vlm_description,
    )

    # Persist the structured signals so re-summarization is idempotent
    # without re-running every stage from scratch and so C9 multi-axis
    # filtering can query them later. Empty / None inputs are dropped.
    #
    # Stash on `image.__dict__` directly (NOT `image.summary_signals = ...`)
    # to avoid going through SQLAlchemy's InstanceState. A setattr on the
    # session-tracked ORM attribute can trigger a state-load if the
    # column's expired in this thread, which calls await_only() without
    # a greenlet and poisons the session for `_mark_done`. Reading from
    # `image.__dict__` in the async wrapper is symmetric and equally safe.
    signals: dict = {}
    if regions: signals["regions"] = regions
    if objects: signals["objects"] = objects
    if concepts: signals["concepts"] = concepts
    if vlm_description: signals["vlm"] = vlm_description
    if signals:
        image.__dict__["summary_signals"] = signals

    if not summary:
        # Deterministic fallback — same path as v1.
        cleaned = _clean_caption(raw_caption)
        spliced = _splice_names(cleaned, named_people)
        summary = spliced or _fallback_image_summary(image, named_people)
        if ocr_text and "text" not in summary.lower():
            excerpt = ocr_text[:120].replace("\n", " ").strip()
            if excerpt:
                summary = f"{summary.rstrip('.')}. Visible text: {excerpt}."

    # Prefer the LLM topic when available — much richer than the heuristic
    # "Photo of Me" fallback. Pass the same signals the rewriter saw so
    # the topic stays consistent with the description.
    topic = _llm_compose_topic(
        caption=raw_caption,
        names=named_people,
        regions=regions,
        objects=objects,
        concepts=concepts,
        scene=image.scene_label,
        setting=image.indoor_outdoor,
        content_type=image.content_type,
    ) or _compose_topic(image, summary, named_people)

    points: list[str] = []
    if named_people:
        # Lead with names — the most useful search anchor.
        points.append("People: " + ", ".join(named_people))
    elif image.face_likelihood and image.face_likelihood > 0.5:
        points.append("Likely contains people")
    if image.content_type and image.content_type != "photo":
        points.append(f"Type: {image.content_type}")
    if image.scene_confidence and image.scene_label:
        points.append(
            f"Scene: {image.scene_label.replace('_', ' ')} "
            f"({int(image.scene_confidence * 100)}%)"
        )
    if image.indoor_outdoor and image.indoor_outdoor != "unknown":
        points.append(f"Setting: {image.indoor_outdoor}")
    if ocr_text:
        excerpt = ocr_text[:80].replace("\n", " ").strip()
        if excerpt:
            ellipsis = "…" if len(ocr_text) > 80 else ""
            points.append(f"Text: {excerpt}{ellipsis}")
    # Top tags from image_tags relationship if loaded; otherwise skip silently.
    tag_labels = []
    try:
        for it in image.tags or []:
            if it.tag is not None and it.tag.label:
                tag_labels.append(it.tag.label)
        if tag_labels:
            points.append("Tags: " + ", ".join(tag_labels[:5]))
    except Exception:
        pass

    return SummaryResult(topic=topic, summary=summary, points=points[:5])


# Filler phrases BLIP frequently leads with — "This is a photo of a man"
# is filler around the actual content "a man". Stripping them keeps the
# searchable terms (subject, scene, objects) and drops the framing.
# Order matters: longer, more specific prefixes must come before shorter
# subsets so we don't half-strip them.
_FILLER_PREFIXES = (
    r"^the (?:image|photo|picture) shows\s+",
    r"^in the (?:image|photo|picture),?\s+",
    r"^in this (?:image|photo|picture),?\s+",
    r"^this is a (?:photo|picture|photograph|image) of\s+",
    r"^this is an? (?:photo|picture|photograph|image)\s+",
    r"^an? (?:image|picture|photo|photograph) of\s+",
    r"^this is\s+",
    r"^there are\s+",
    r"^there is\s+",
)


# Patterns BLIP uses for generic person references, in order of priority
# (most specific first). Each gets exactly one substitution per caption.
_PERSON_PATTERNS = (
    r"\btwo people\b",
    r"\bthree people\b",
    r"\bseveral people\b",
    r"\ba group of people\b",
    r"\bgroup of people\b",
    r"\bpeople\b",
    r"\ba young man\b",
    r"\ba young woman\b",
    r"\ba man\b",
    r"\ba woman\b",
    r"\ba boy\b",
    r"\ba girl\b",
    r"\ba person\b",
    r"\bthe man\b",
    r"\bthe woman\b",
)


def _clean_caption(caption: str) -> str:
    """Strip BLIP filler and normalize awkward phrasings before name splicing.

    BLIP often pads captions with "This is a..." / "There is a..." which
    adds no semantic value and dilutes search relevance — the user is
    searching for "whiteboard math classroom", not for "this is a photo".
    Removing the filler also lets `_splice_names` produce natural sentences:
    "There is a man writing" → "There is Mr Koler writing" (awkward) becomes
    "A man writing" → "Mr Koler writing" (clean).

    Also rewrites "taking a picture/photo of himself|herself|themselves" to
    "taking a selfie", which is the term people actually search by.
    """
    if not caption:
        return caption

    cleaned = caption
    for pattern in _FILLER_PREFIXES:
        new, n = re.subn(pattern, "", cleaned, count=1, flags=re.IGNORECASE)
        if n > 0:
            cleaned = new
            break  # one prefix at most — chained strips would over-cut

    # Plural first so "pictures of themselves" doesn't get half-matched by
    # the singular pattern.
    cleaned = re.sub(
        r"\btaking (?:pictures|photos) of themselves\b",
        "taking selfies",
        cleaned,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(
        r"\btaking a (?:picture|photo) of (?:himself|herself|themselves)\b",
        "taking a selfie",
        cleaned,
        flags=re.IGNORECASE,
    )

    cleaned = cleaned.strip()
    if cleaned:
        cleaned = cleaned[0].upper() + cleaned[1:]
        if cleaned[-1] not in ".!?":
            cleaned += "."
    return cleaned


def _splice_names(caption: str, names: list[str]) -> str:
    """Substitute the first generic person reference with the named people.

    "A black and white photo of a man looking out a car window."
        + names=["Me"]
    → "A black and white photo of Me looking out a car window."

    Returns the caption unchanged when nothing matches or no names are
    known, so BLIP captions stay grammatical and search-friendly.
    """
    if not caption or not names:
        return caption
    if len(names) == 1:
        replacement = names[0]
    elif len(names) == 2:
        replacement = f"{names[0]} and {names[1]}"
    else:
        replacement = ", ".join(names[:-1]) + f", and {names[-1]}"

    spliced = caption
    matched = False
    for pattern in _PERSON_PATTERNS:
        new_caption, n = re.subn(
            pattern, replacement, spliced, count=1, flags=re.IGNORECASE
        )
        if n > 0:
            spliced = new_caption
            matched = True
            break
    if not matched:
        # No person pattern matched (e.g. landscape) — append parenthetical
        # so names still index for search without mangling the sentence.
        return f"{caption} (with {replacement})"

    return _polish_after_splice(spliced, names)


def _splice_names_all(caption: str, names: list[str]) -> str:
    """Replace every generic person reference in a (possibly multi-sentence)
    caption with the named people. Sister to `_splice_names` for the
    Qwen rewriter output path.

    The single-replacement `_splice_names` works for BLIP's one-clause
    captions ("a man taking a selfie" → "Me taking a selfie"). The C2
    Qwen rewriter emits 1–3 sentence descriptions that mention the
    subject multiple times ("A young man stands… The man wears… He
    holds…"). Replacing only the first reference left every later
    "the man" / "he" / "his" pointing at an unnamed entity, which is
    exactly the "summaries don't replace pronouns with detected
    people" complaint.

    Strategy:
      1. Replace each generic-noun pattern (a/the man, young man,
         person, etc.) with the name — but only the *first*
         occurrence per pattern, so we don't double-replace already-
         normalized references in the same sentence.
      2. Run the existing first-person pronoun polish so "his/her/him"
         flip to "my/me" when the subject is "Me" / "I".

    Returns the caption unchanged when `names` is empty.
    """
    if not caption or not names:
        return caption
    if len(names) == 1:
        replacement = names[0]
    elif len(names) == 2:
        replacement = f"{names[0]} and {names[1]}"
    else:
        replacement = ", ".join(names[:-1]) + f", and {names[-1]}"

    spliced = caption
    # First substitution per pattern, scanning each pattern in turn —
    # avoids cascading replacements like "a young man" → "Me" then
    # "a man" matching part of the new text. Patterns are ordered
    # specific-to-general in `_PERSON_PATTERNS` already.
    for pattern in _PERSON_PATTERNS:
        spliced, _ = re.subn(
            pattern, replacement, spliced, count=0, flags=re.IGNORECASE
        )

    return _polish_after_splice(spliced, names)


# Pronoun / grammar cleanup that only makes sense once names are inserted.
# BLIP says "a man that is taking a picture of himself with his phone";
# splicing yields "Me that is taking a picture of himself with his phone";
# this pass rewrites it to "Me taking a selfie with my phone".
_FIRST_PERSON_NAMES = {"me", "i"}


def _polish_after_splice(caption: str, names: list[str]) -> str:
    # 1. "[name] that is V-ing" / "[name] who is V-ing" → "[name] V-ing".
    #    Reads natural in English ("Mr Koler standing" beats "Mr Koler that
    #    is standing") and is what BLIP really meant.
    for name in names:
        esc = re.escape(name)
        caption = re.sub(
            rf"\b{esc}\s+(?:that|who)\s+is\s+",
            f"{name} ",
            caption,
        )

    # 2. First-person pronoun rewrite. When the spliced subject is "Me" or
    #    "I", any third-person pronouns BLIP emitted refer back to that
    #    subject and need to flip to first person — otherwise we get
    #    "Me taking a selfie with his cell phone".
    if any(n.lower() in _FIRST_PERSON_NAMES for n in names):
        caption = re.sub(
            r"\b(?:himself|herself|themselves)\b",
            "myself",
            caption,
            flags=re.IGNORECASE,
        )
        # "his" is unambiguously possessive; "her" is mostly possessive in
        # caption contexts ("her phone", "her face"), so map both to "my".
        caption = re.sub(r"\bhis\b", "my", caption, flags=re.IGNORECASE)
        caption = re.sub(r"\bher\b", "my", caption, flags=re.IGNORECASE)
        # Object pronoun "him" is rare but happens ("a dog next to him").
        caption = re.sub(r"\bhim\b", "me", caption, flags=re.IGNORECASE)

    return caption


def _compose_topic(
    image: Image, caption: str, named_people: list[str]
) -> str:
    """Pick a short, content-rich topic line. Named people lead — Drive's
    "People" facet works the same way."""
    cap_lower = caption.lower()
    is_selfie = "selfie" in cap_lower or "self portrait" in cap_lower
    is_group = (
        "group photo" in cap_lower
        or "group of people" in cap_lower
        or len(named_people) >= 2
    )

    if named_people:
        names = ", ".join(named_people)
        if is_selfie:
            return f"Selfie of {names}"
        if is_group:
            return f"Group photo · {names}"
        return f"Photo of {names}"

    if is_selfie:
        return "Selfie"
    if is_group:
        return "Group photo"
    if image.scene_label == "portrait" or image.scene_label == "group_photo":
        return image.scene_label.replace("_", " ").title()
    if image.face_likelihood and image.face_likelihood > 0.7 and (
        "person" in cap_lower or "man" in cap_lower or "woman" in cap_lower
    ):
        return "Portrait"
    if image.scene_label and image.scene_label not in ("portrait",):
        return image.scene_label.replace("_", " ").title()
    return "Photo"


def _fallback_image_summary(image: Image, named_people: list[str] | None = None) -> str:
    """Last-resort summary built from classifier columns + identified people.

    The earlier version produced "A classroom indoor image." with no
    reference to the named subjects — defeating the point of running the
    face pipeline in the first place. When we do have names, lead with
    them so the description reads "Photo of Mr Koler in a classroom
    (indoor)." which is both more useful and matches what a user would
    actually search by.
    """
    scene = (image.scene_label or "").replace("_", " ").strip()
    indoor = image.indoor_outdoor if image.indoor_outdoor and image.indoor_outdoor != "unknown" else None
    names = [n for n in (named_people or []) if n]
    if names:
        if len(names) == 1:
            who = names[0]
        elif len(names) == 2:
            who = f"{names[0]} and {names[1]}"
        else:
            who = ", ".join(names[:-1]) + f", and {names[-1]}"
        tail = []
        if scene:
            tail.append(f"in a {scene}" if scene[0].lower() not in "aeiou" else f"in an {scene}")
        if indoor:
            tail.append(f"({indoor})")
        return f"Photo of {who} " + " ".join(tail) + "." if tail else f"Photo of {who}."
    bits = []
    if scene:
        bits.append(scene)
    if indoor:
        bits.append(indoor)
    if not bits:
        return "Image."
    return "A " + " ".join(bits) + " image."


def _caption_image(raw_bytes: bytes) -> Optional[str]:
    """Detailed image caption. Florence-2 primary, BLIP fallback.

    Returns a 1-3 sentence detailed description, or None if both models
    fail to load. The caller folds None into the regex/fallback path.
    """
    cap = _florence_caption(raw_bytes)
    if cap:
        return cap
    return _blip_caption(raw_bytes)


def _florence_caption(raw_bytes: bytes) -> Optional[str]:
    """Florence-2 <MORE_DETAILED_CAPTION>. Returns None on any failure.

    Outputs are dense ("The image shows a man wearing glasses standing in
    front of a large whiteboard covered in handwritten mathematical
    equations and diagrams in blue marker"), much richer than BLIP's
    one-clause captions. Filler ("The image shows...") is stripped by the
    LLM rewriter or by `_clean_caption` in the regex fallback path.
    """
    global _FLORENCE_BROKEN
    if _FLORENCE_BROKEN:
        return None
    try:
        from PIL import Image as PILImage
        import torch

        from backend.vision.runtime import get_florence2

        model, processor, device = get_florence2()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        prompt = "<MORE_DETAILED_CAPTION>"
        inputs = processor(text=prompt, images=image, return_tensors="pt")
        inputs = {k: v.to(device) for k, v in inputs.items()}
        if device == "cuda" and "pixel_values" in inputs:
            inputs["pixel_values"] = inputs["pixel_values"].half()

        with torch.no_grad():
            ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=256,
                num_beams=3,
                do_sample=False,
            )
        text = processor.batch_decode(ids, skip_special_tokens=False)[0]
        parsed = processor.post_process_generation(
            text, task=prompt, image_size=(image.width, image.height)
        )
        out = (parsed.get(prompt) or "").strip()
        return out or None
    except KeyError as e:
        # transformers >=4.50 Cache-shape mismatch — Florence-2's
        # bundled modeling code is incompatible. Disable globally so
        # BLIP fallback runs instantly on every later image.
        logger.warning("florence2: caption disabled (KeyError: %s)", e)
        _FLORENCE_BROKEN = True
        return None
    except Exception:
        logger.exception("florence2: caption failed")
        return None


def _ocr_image(raw_bytes: bytes) -> Optional[str]:
    """Florence-2 <OCR>. Returns extracted text or None.

    Reads visible text in the image — whiteboards, document scans,
    screenshots. Empty/whitespace-only output → None so callers don't
    propagate a useless empty string.
    """
    global _FLORENCE_BROKEN
    if _FLORENCE_BROKEN:
        return None
    try:
        from PIL import Image as PILImage
        import torch

        from backend.vision.runtime import get_florence2

        model, processor, device = get_florence2()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        prompt = "<OCR>"
        inputs = processor(text=prompt, images=image, return_tensors="pt")
        inputs = {k: v.to(device) for k, v in inputs.items()}
        if device == "cuda" and "pixel_values" in inputs:
            inputs["pixel_values"] = inputs["pixel_values"].half()

        with torch.no_grad():
            ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=512,
                num_beams=3,
                do_sample=False,
            )
        text = processor.batch_decode(ids, skip_special_tokens=False)[0]
        parsed = processor.post_process_generation(
            text, task=prompt, image_size=(image.width, image.height)
        )
        ocr = (parsed.get(prompt) or "").strip()
        return ocr or None
    except KeyError as e:
        logger.warning("florence2: ocr disabled (KeyError: %s)", e)
        _FLORENCE_BROKEN = True
        return None
    except Exception:
        logger.exception("florence2: ocr failed")
        return None


# Florence-2's bundled `prepare_inputs_for_generation` expects the
# pre-Cache-object past_key_values shape (tuple of tuples). With
# transformers >= 4.50 that shape is gone and EVERY Florence-2 task
# (caption, OCR, dense regions, OD) raises
# `KeyError: Cache only has 0 layers` the moment beam search starts.
# The bundled code doesn't get upgraded in lock-step, so once we see
# the failure we mark Florence-2 entirely disabled — every later
# image falls straight through to BLIP for captioning + Qwen for
# rewriting, instead of wasting ~30 s per task reloading the model
# and re-crashing in the same spot.
_FLORENCE_BROKEN: bool = False


def _florence_regions(raw_bytes: bytes) -> Optional[list[str]]:
    """Florence-2 `<DENSE_REGION_CAPTION>` — per-region descriptive
    phrases like "a kettle on the floor" or "a row of mathematical
    equations in blue marker."

    Returns a deduped list of phrases (typically 5–25 per image) or
    None on failure. Feeds the Qwen synthesis prompt as
    "Objects and regions:" context so the final summary mentions
    grounded specifics instead of generic captions.
    """
    global _FLORENCE_BROKEN
    if _FLORENCE_BROKEN:
        return None
    try:
        from PIL import Image as PILImage
        import torch

        from backend.vision.runtime import get_florence2

        model, processor, device = get_florence2()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        prompt = "<DENSE_REGION_CAPTION>"
        inputs = processor(text=prompt, images=image, return_tensors="pt")
        inputs = {k: v.to(device) for k, v in inputs.items()}
        if device == "cuda" and "pixel_values" in inputs:
            inputs["pixel_values"] = inputs["pixel_values"].half()

        with torch.no_grad():
            ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=1024,
                num_beams=3,
                do_sample=False,
            )
        text = processor.batch_decode(ids, skip_special_tokens=False)[0]
        parsed = processor.post_process_generation(
            text, task=prompt, image_size=(image.width, image.height)
        )
        # Florence-2 DENSE_REGION_CAPTION returns either
        # {"labels": [...], "bboxes": [...]} or a similar shape — the exact
        # keys vary across transformers versions. Handle both.
        result = parsed.get(prompt) or {}
        labels = result.get("labels") or result.get("captions") or []
        # Dedup while preserving order, strip whitespace, drop empties.
        seen = set()
        out: list[str] = []
        for label in labels:
            s = (label or "").strip()
            if not s or s.lower() in seen:
                continue
            seen.add(s.lower())
            out.append(s)
        return out or None
    except KeyError as e:
        # Cache shape mismatch from transformers >= 4.50 — flip the
        # permanent disable flag so we stop re-trying every image.
        logger.warning("florence2: dense region caption disabled (KeyError: %s)", e)
        _FLORENCE_BROKEN = True
        return None
    except Exception:
        logger.exception("florence2: dense region caption failed")
        return None


def _florence_objects(raw_bytes: bytes) -> Optional[list[str]]:
    """Florence-2 `<OD>` (object detection with labels). Returns a
    deduped list of detected object labels (e.g. "person", "laptop",
    "coffee mug"). Independent from `_florence_regions`: this gives
    canonical category names useful for filtering / search, while
    regions give descriptive phrases for the LLM prompt.
    """
    global _FLORENCE_BROKEN
    if _FLORENCE_BROKEN:
        return None
    try:
        from PIL import Image as PILImage
        import torch

        from backend.vision.runtime import get_florence2

        model, processor, device = get_florence2()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        prompt = "<OD>"
        inputs = processor(text=prompt, images=image, return_tensors="pt")
        inputs = {k: v.to(device) for k, v in inputs.items()}
        if device == "cuda" and "pixel_values" in inputs:
            inputs["pixel_values"] = inputs["pixel_values"].half()

        with torch.no_grad():
            ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=1024,
                num_beams=3,
                do_sample=False,
            )
        text = processor.batch_decode(ids, skip_special_tokens=False)[0]
        parsed = processor.post_process_generation(
            text, task=prompt, image_size=(image.width, image.height)
        )
        result = parsed.get(prompt) or {}
        labels = result.get("labels") or []
        seen = set()
        out: list[str] = []
        for label in labels:
            s = (label or "").strip().lower()
            if not s or s in seen:
                continue
            seen.add(s)
            out.append(s)
        return out or None
    except KeyError as e:
        logger.warning("florence2: OD disabled (KeyError: %s)", e)
        _FLORENCE_BROKEN = True
        return None
    except Exception:
        logger.exception("florence2: OD failed")
        return None


def _vlm_describe(raw_bytes: bytes) -> Optional[str]:
    """C2e — heavy vision-language description via InternVL2-4B.

    Returns a 1-paragraph natural-language description of the image
    that captures nuance Florence-2 misses (mood, action, fine-grained
    object relationships). Gated behind `settings.heavy_vlm_enabled`;
    returns None whenever the model can't load — the rest of the
    summary pipeline degrades cleanly.

    InternVL2's chat API expects a `<image>` token in the prompt that
    the model's `chat()` method swaps for image embeddings. Exact
    invocation varies between versions; we use the canonical pattern
    from the InternVL2 model card and catch any incompatibility.
    """
    try:
        import torch  # type: ignore
        from PIL import Image as PILImage  # type: ignore

        from backend.vision.runtime import get_internvl2

        model, tokenizer, device = get_internvl2()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        # InternVL2 exposes a `chat()` helper on the model that handles
        # image preprocessing internally given a PIL image. The exact
        # signature is `model.chat(tokenizer, pixel_values, question,
        # generation_config, history=None)` where pixel_values comes
        # from the model's `_preprocess_image` (also exposed as a class
        # method on the modeling file). We use the simpler `chat` path
        # that accepts a PIL image directly — newer revisions support it.
        question = (
            "Describe this image in 2-3 specific sentences. Cover every "
            "visible person, object, action, environment cue, and any "
            "readable text. Use concrete nouns; avoid filler like 'The "
            "image shows'. Output only the description."
        )
        generation_config = dict(
            max_new_tokens=256,
            do_sample=False,
        )

        # Run inference. Wrap broadly — different InternVL2 builds expose
        # slightly different surfaces.
        with torch.no_grad():
            if hasattr(model, "chat"):
                # Newer builds: model.chat(tokenizer, image, question, ...)
                try:
                    out = model.chat(
                        tokenizer,
                        image,
                        question,
                        generation_config=generation_config,
                    )
                except TypeError:
                    out = None
                if isinstance(out, str):
                    text = out.strip()
                elif isinstance(out, tuple) and out and isinstance(out[0], str):
                    text = out[0].strip()
                else:
                    text = None
            else:
                text = None

        if not text:
            return None
        # Collapse whitespace; cap at 1000 chars so a runaway model
        # doesn't dump a wall of text into the Qwen prompt budget.
        text = re.sub(r"\s+", " ", text).strip().strip('"').strip("'").strip()
        if len(text) > 1000:
            text = text[:1000].rsplit(".", 1)[0] + "."
        return text or None
    except Exception:
        logger.exception("vlm: heavy VLM description failed")
        return None


def _blip_caption(raw_bytes: bytes) -> Optional[str]:
    """BLIP fallback when Florence-2 fails. Same shape as v1.

    BLIP produces short single-clause captions ("a man taking a selfie in
    a kitchen"); kept as a safety net so first-time installs that haven't
    pulled Florence-2 weights still work.
    """
    try:
        from PIL import Image as PILImage
        import torch

        from backend.vision.runtime import get_caption_model

        model, processor, device = get_caption_model()
        image = PILImage.open(BytesIO(raw_bytes)).convert("RGB")

        inputs = processor(images=image, return_tensors="pt")
        inputs = {k: v.to(device) for k, v in inputs.items()}
        if device == "cuda" and "pixel_values" in inputs:
            inputs["pixel_values"] = inputs["pixel_values"].half()

        with torch.no_grad():
            generated_ids = model.generate(
                **inputs,
                max_new_tokens=80,
                num_beams=5,
                do_sample=False,
                repetition_penalty=1.2,
            )
        caption = processor.decode(generated_ids[0], skip_special_tokens=True)
        caption = caption.strip()
        if caption:
            caption = caption[0].upper() + caption[1:]
            if caption[-1] not in ".!?":
                caption += "."
        return caption or None
    except Exception:
        logger.exception("blip: caption failed")
        return None


def _llm_rewrite_summary(
    caption: str,
    names: list[str],
    ocr_text: Optional[str],
    scene: Optional[str],
    setting: Optional[str],
    content_type: Optional[str],
    tags: Optional[list[str]] = None,
    regions: Optional[list[str]] = None,
    objects: Optional[list[str]] = None,
    concepts: Optional[list[str]] = None,
    vlm_description: Optional[str] = None,
) -> Optional[str]:
    """Qwen2.5-Instruct synthesizes raw multi-model signals into one
    grounded description.

    Accepts everything the C2 pipeline produces:
      - `caption`: Florence-2 <MORE_DETAILED_CAPTION>.
      - `ocr_text`: Florence-2 <OCR> (scene-gated).
      - `regions`: Florence-2 <DENSE_REGION_CAPTION> phrases.
      - `objects`: Florence-2 <OD> label list.
      - `concepts`: OpenCLIP top-K against the curated concept vocab.
      - `vlm_description`: optional InternVL2-4B paragraph (heavy, gated).
      - `tags`: any existing Image.tags labels for the row.
      - plus `scene`, `setting`, `content_type` from the existing
        classifier pass.

    Each signal is optional — missing inputs just drop the corresponding
    context line, so a stage that errors out degrades the summary
    instead of breaking it. Returns None on rewriter failure; caller
    falls back to the regex pipeline.
    """
    if not settings.rewriter_enabled:
        return None
    if not caption and not ocr_text and not vlm_description and not regions:
        return None

    try:
        import torch

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        ctx_lines: list[str] = []
        if caption:
            ctx_lines.append(f"Caption: {caption}")
        if vlm_description:
            ctx_lines.append(f"Rich description: {vlm_description}")
        if names:
            ctx_lines.append(f"People in image: {', '.join(names)}")
        if regions:
            ctx_lines.append(
                "Objects and regions:\n- " + "\n- ".join(regions[:30])
            )
        if objects:
            ctx_lines.append(f"Detected objects: {', '.join(objects[:30])}")
        if concepts:
            ctx_lines.append(f"Concept tags: {', '.join(concepts[:20])}")
        if tags:
            ctx_lines.append(f"Existing tags: {', '.join(tags[:20])}")
        if ocr_text:
            # `ocr_text` is repurposed by the video summarizer to carry
            # the audio transcript (see `_summarize_video`). Label
            # accordingly so the LLM knows whether it's looking at
            # whiteboard OCR (image path) or spoken content (video
            # path). 1500 chars fits the prompt budget alongside the
            # multi-signal context above.
            label = (
                "Spoken content (what was said) — this carries the "
                "actual subject of the video"
                if content_type == "video"
                else "Visible text in image"
            )
            ctx_lines.append(f"{label}: {ocr_text[:1500]}")
        if scene:
            ctx_lines.append(f"Scene: {scene.replace('_', ' ')}")
        if setting and setting != "unknown":
            ctx_lines.append(f"Setting: {setting}")
        if content_type:
            ctx_lines.append(f"Content type: {content_type}")

        first_person = any(
            (n or "").strip().lower() in {"me", "i"} for n in names
        )

        # Video and image summaries need DIFFERENT prompts. The image
        # path emphasises visual nouns and treats on-screen text as
        # "describe what it's ABOUT" — for video that wipes out the
        # actual spoken content, which the user cares about MORE than
        # the visual frame. The video path treats transcript +
        # on-screen text as half the signal (the SUBJECT of the
        # clip) and the keyframe visuals as the other half.
        if content_type == "video":
            # Prompt rewrite (2026-05): the previous version had "(a) "
            # and "(b) WHAT WAS SAID" + "what the camera shows" as
            # numbered instruction clauses. Weaker LLM checkpoints
            # interpreted those as STRUCTURAL HEADERS and emitted
            # output like "**WHAT WAS SAID:** … **ON-SCREEN TEXT:** …
            # **VISUAL SCENE:** …" — section markdown that leaked
            # back to the user. The new instructions describe the
            # blending intent without giving the model a literal
            # template to copy.
            #
            # Also: paraphrase MORE aggressively. The previous version
            # ended up producing near-verbatim transcripts when the
            # transcript was short and Florence captions thin —
            # Qwen took the path of least resistance and just echoed
            # the spoken text. We now explicitly forbid quoting.
            instructions = (
                "Write a single paragraph (2-4 sentences, up to ~90 "
                "words) describing this video as ONE coherent whole. "
                "Blend what was said in the audio with what's visible "
                "on screen — produce a single narrative, NOT a "
                "frame-by-frame breakdown. Do NOT mention 'the first "
                "frame', 'the second screenshot', or any other "
                "enumeration of the keyframes — synthesize across "
                "them. Do NOT produce sections, headings, bullet "
                "points, or labels like 'WHAT WAS SAID'. Paraphrase "
                "the spoken content into your own words; do not "
                "quote it verbatim or repeat full sentences from the "
                "transcript. Name concrete topics, people, places, "
                "and decisions discussed instead of saying 'they talk "
                "about technology'. Mention the visual setting "
                "briefly. Use named people instead of 'a man'. Do "
                "NOT begin with 'The video shows', 'This is a', or "
                "'There is'. Output only the description."
            )
        else:
            instructions = (
                "Write a dense, keyword-rich description (1-3 sentences, up "
                "to ~70 words) of what's in this image. Pack in EVERY "
                "distinct concrete noun, named person, scene cue, lighting "
                "detail, action, and object from the inputs above — these "
                "are the search keywords users will type later, so "
                "redundancy with the inputs is good, not bad. "
                "Use named people instead of phrases like 'a man'. "
                "If visible text is provided, describe what the text is "
                "ABOUT (e.g. 'matrix algebra equations', 'a chat "
                "conversation', 'a recipe') rather than quoting it verbatim. "
                "Do NOT begin with 'The image shows', 'This is a', or "
                "'There is'. Output only the description — no preamble, no "
                "quotes, no bullet lists."
            )
        if first_person:
            instructions += (
                " The named person 'Me' is the photo owner — use first "
                "person ('I', 'my', 'myself')."
            )

        messages = [
            {
                "role": "system",
                "content": (
                    "You write natural, search-friendly descriptions from "
                    "structured image-analysis signals. You are concrete, "
                    "factual, and never invent details that aren't in the "
                    "input."
                ),
            },
            {
                "role": "user",
                "content": instructions + "\n\n" + "\n".join(ctx_lines),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(prompt, return_tensors="pt").to(device)
        prompt_len = inputs.input_ids.shape[1]

        # Video summaries are 2-4 sentences (~90 words = ~120 tokens
        # for English). Image summaries are 1-3 sentences (~70
        # words = ~95 tokens). Give video a bigger budget so the
        # generation doesn't get cut mid-sentence and produce
        # something that fails the rejection cap.
        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=220 if content_type == "video" else 140,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()

        # Guard against the LLM going off — JSON dumps, multi-paragraph
        # ramble, etc. Collapse internal whitespace.
        # Caps: 700 chars for image (1-3 sentence) summaries, 1300
        # chars for video (2-4 sentence, ~90 words) summaries. The
        # old uniform 700-char cap rejected legitimate 90-word video
        # outputs and fell back to the raw transcript verbatim.
        reply = re.sub(r"\s+", " ", reply).strip().strip('"').strip("'").strip()
        # Strip leaked markdown section headers — Qwen sometimes emits
        # **WHAT WAS SAID:** / **VISUAL SCENE:** style prefixes despite
        # the instructions telling it not to. Drop them so the user
        # sees clean prose.
        if content_type == "video":
            reply = re.sub(
                r"\*\*[A-Z][A-Z0-9 _\-]+\*\*[:\s]*", " ", reply,
            ).strip()
            # Strip frame-by-frame enumeration. Qwen-2.5-1.5B is a
            # small model and treats the multi-keyframe Florence
            # caption list as a per-frame instruction set despite
            # the prompt forbidding it — outputs like "In the first
            # frame ... In the second frame ..." dominate when
            # there isn't much spoken content. We rewrite the most
            # common patterns into neutral connective tissue so
            # the summary still parses as one paragraph instead of
            # a numbered breakdown.
            FRAME_ENUM = re.compile(
                r"\b(?:in\s+|at\s+)?(?:the\s+)?"
                r"(?:first|second|third|fourth|fifth|sixth|seventh|"
                r"eighth|next|final|last|opening|closing|"
                r"earlier|later|following|preceding)"
                r"\s+(?:frame|screenshot|image|shot|scene|clip|"
                r"keyframe|snapshot|moment),?\s*",
                re.IGNORECASE,
            )
            reply = FRAME_ENUM.sub("", reply)
            # Also drop the conjunctive "Each frame captures a
            # unique aspect" / "Each screenshot shows" boilerplate
            # — same model tic, lower frequency.
            reply = re.sub(
                r"\b(?:Each|Every)\s+(?:frame|screenshot|image|shot|"
                r"keyframe|snapshot)\s+(?:captures|shows|features|"
                r"depicts|displays|reveals)\b[^.]*\.\s*",
                "",
                reply, flags=re.IGNORECASE,
            )
            # Collapse the gaps the stripping left behind. Capitalize
            # the start of the first surviving sentence.
            reply = re.sub(r"\s+", " ", reply).strip()
            reply = re.sub(r"\.\s*([a-z])", lambda m: ". " + m.group(1).upper(), reply)
            if reply and reply[0].islower():
                reply = reply[0].upper() + reply[1:]
        max_len = 1300 if content_type == "video" else 700
        if not reply or len(reply) > max_len:
            return None
        # Ensure terminal punctuation.
        if reply[-1] not in ".!?":
            reply = reply.rstrip(",;:") + "."
        # Capitalize.
        reply = reply[0].upper() + reply[1:]
        # Post-process: substitute every generic person reference
        # ("a man", "the man", "a young man", "he", etc.) with the
        # detected named people. Qwen2.5-1.5B follows the "Use named
        # people" instruction only sometimes — its multi-sentence
        # descriptions often re-introduce "the man" / "he" mid-paragraph
        # because each sentence is generated against the prior caption,
        # not the people-list metadata. Running the regex splice on the
        # final output makes the substitution deterministic so summaries
        # are always indexed under the user's actual search terms.
        if names:
            reply = _splice_names_all(reply, names)
        return reply
    except Exception:
        logger.exception("rewriter: failed")
        return None


def _llm_compose_topic(
    caption: Optional[str],
    names: list[str],
    regions: Optional[list[str]] = None,
    objects: Optional[list[str]] = None,
    concepts: Optional[list[str]] = None,
    scene: Optional[str] = None,
    setting: Optional[str] = None,
    content_type: Optional[str] = None,
) -> Optional[str]:
    """Ask Qwen for a short search-friendly topic line (6-12 words).

    Replaces the v1 heuristic that produced "Photo of Me" / "Selfie" /
    "Photo of Mr Koler" — too generic for search. The topic should
    surface 1-2 specific nouns plus the named person when applicable
    ("Selfie at home in a desert-themed room", "Mr Koler at a
    whiteboard with linear algebra equations").

    Same context signals as the rewriter so topic + summary stay
    consistent. Returns None on failure; caller falls back to the
    heuristic `_compose_topic`.
    """
    if not settings.rewriter_enabled:
        return None
    if not caption and not regions and not objects:
        return None

    try:
        import torch  # type: ignore

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        ctx_lines: list[str] = []
        if caption:
            ctx_lines.append(f"Caption: {caption}")
        if names:
            ctx_lines.append(f"People: {', '.join(names)}")
        if regions:
            ctx_lines.append(f"Regions: {'; '.join(regions[:10])}")
        if objects:
            ctx_lines.append(f"Objects: {', '.join(objects[:15])}")
        if concepts:
            ctx_lines.append(f"Concepts: {', '.join(concepts[:10])}")
        if scene:
            ctx_lines.append(f"Scene: {scene.replace('_', ' ')}")
        if setting and setting != "unknown":
            ctx_lines.append(f"Setting: {setting}")
        if content_type:
            ctx_lines.append(f"Content type: {content_type}")

        instructions = (
            "Write a 6-12 word topic line for this image that a user "
            "could type into a search bar to find it again. Include "
            "the most distinctive 1-3 specific nouns or names from "
            "the inputs (a person's name, the place, an activity, a "
            "specific object). No filler words. Title-case is fine but "
            "not required. No trailing punctuation. Output only the "
            "topic line — no quotes, no preamble."
        )

        messages = [
            {
                "role": "system",
                "content": (
                    "You write short search-friendly topic lines from "
                    "structured image-analysis signals. Concrete, "
                    "specific, never inventive."
                ),
            },
            {
                "role": "user",
                "content": instructions + "\n\n" + "\n".join(ctx_lines),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(prompt, return_tensors="pt").to(device)
        prompt_len = inputs.input_ids.shape[1]

        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=40,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        reply = re.sub(r"\s+", " ", reply).strip()
        # Strip surrounding quotes, leading/trailing punctuation,
        # 'Topic:' / 'Title:' prefixes the LLM occasionally adds.
        reply = reply.strip('"').strip("'").strip()
        reply = re.sub(r"^(topic|title)\s*:\s*", "", reply, flags=re.IGNORECASE)
        # Drop the trailing period some models add.
        reply = reply.rstrip(".,;:!? ")
        if not reply or len(reply) > 100:
            return None
        # Take only the first line if the model returned multiple.
        reply = reply.split("\n")[0].strip()
        if not reply:
            return None
        return reply
    except Exception:
        logger.exception("topic rewriter: failed")
        return None


def _llm_compose_doc_topic(text: str, filename: str) -> Optional[str]:
    """Same as `_llm_compose_topic` but for documents. Caller passes the
    extracted text and original filename; we produce a 6-12 word topic
    line that surfaces the most distinctive nouns / sections / numbers.
    Returns None on failure; caller falls back to the heuristic
    (filename stem)."""
    if not settings.rewriter_enabled:
        return None
    if not text:
        return None

    try:
        import torch  # type: ignore

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        # Keep the doc head short — topic lines don't need the whole
        # body and a smaller context speeds up the call.
        head = text[:2000]
        instructions = (
            "Write a 6-12 word topic line for this document that a "
            "user could type into a search bar to find it again. "
            "Use the most distinctive concrete nouns from the content "
            "— project names, sections, numbers, dates, parties. "
            "No filler words. No trailing punctuation. Output only "
            "the topic — no quotes, no preamble, no 'Title:' prefix."
        )

        messages = [
            {
                "role": "system",
                "content": (
                    "You write short search-friendly topic lines for "
                    "documents. Concrete, specific, never inventive."
                ),
            },
            {
                "role": "user",
                "content": (
                    instructions
                    + f"\n\nFilename: {filename}\n\nContent:\n{head}"
                ),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3500
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]

        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=40,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        reply = re.sub(r"\s+", " ", reply).strip().strip('"').strip("'").strip()
        reply = re.sub(r"^(topic|title)\s*:\s*", "", reply, flags=re.IGNORECASE)
        reply = reply.rstrip(".,;:!? ")
        if not reply or len(reply) > 100:
            return None
        reply = reply.split("\n")[0].strip()
        return reply or None
    except Exception:
        logger.exception("doc topic rewriter: failed")
        return None


# --- video -----------------------------------------------------------------


def _summarize_video(image: Image, raw_bytes: bytes) -> Optional[SummaryResult]:
    """Multi-keyframe video summary.

    1. Probe duration via ffprobe.
    2. Sample 4 keyframes evenly across the duration (15% / 40% / 65% / 90%
       — avoid the very start/end which are often black or logos).
    3. Florence-caption each frame.
    4. Send the caption set + filename to Qwen via `_llm_rewrite_summary`
       — Qwen synthesizes one coherent description from the per-frame
       observations. Same path image summaries already use.
    5. Topic: short Qwen-style title derived from the aggregated summary.
       Falls back to a cleaned filename when Qwen produces nothing.
    6. Content-type / scene inference: run the existing CLIP concept
       vocabulary over the middle keyframe and surface the top hits as
       `summary_signals.concepts` (search uses these later).

    Each step degrades independently:
      - No ffmpeg → filename-only summary
      - Florence failed on every frame → filename-only summary
      - Qwen disabled / unavailable → concatenated raw captions
    so installs without [ml] extras still get *something* useful.
    """
    duration_s = _probe_video_duration(raw_bytes)
    if duration_s and duration_s > 0:
        # Sampling density scales with duration so a 30s clip gets a
        # tight read (≤6 samples → one every ~5s) and a 30min lecture
        # gets enough coverage to catch scene changes without
        # Florence-captioning every frame at runtime cost.
        # Cap at MAX_FRAMES so a 2h movie doesn't try to caption 720
        # frames (90+ minutes of GPU time).
        #
        # User feedback (2026-05) — bumped sample density from 1/5s to
        # 1/3s. The old rate undersampled short phone clips (a 15-second
        # video got only 4 frames covering 10-90%, missing the
        # information-dense middle). 1 sample / 3s is still cheap
        # on GPU (~50% more Florence calls) and gives meaningful scene
        # coverage on the 10-60s clips that dominate the dataset, while
        # the MAX_FRAMES cap of 30 keeps long videos bounded.
        MAX_FRAMES = 30
        SECONDS_PER_FRAME = 3.0
        n = max(4, min(MAX_FRAMES, int(round(duration_s / SECONDS_PER_FRAME))))
        # Spread evenly inside [10%, 90%] so we skip the standard
        # intro/outro black + fade and concentrate on the actual
        # content. Bias toward earlier coverage where most clips
        # peak in information density.
        if n == 1:
            offsets = [duration_s * 0.5]
        else:
            offsets = [
                max(0.5, duration_s * (0.10 + (0.80 * i / (n - 1))))
                for i in range(n)
            ]
    else:
        # Unknown duration — fall back to fixed offsets, the last of
        # which (60s) will fail silently on shorter clips. ffmpeg just
        # returns the last frame at that point.
        offsets = [1.0, 3.0, 5.0, 10.0, 15.0, 30.0, 60.0]
    frames = [_extract_keyframe(raw_bytes, t) for t in offsets]
    frames = [f for f in frames if f]

    # Sprint I#6 — scene-cut detection. Drop frames whose luminance
    # histogram barely differs from the previous kept frame, so a
    # static talking-head clip captions one frame instead of 30
    # near-identical ones. Keeps Florence cost down + gives Qwen
    # distinct scenes rather than the same observation repeated.
    frames_before_scenecut = len(frames)
    frames = dedup_frames_by_histogram(frames)
    scenes_kept = len(frames)

    # Caption-quality filter — Florence occasionally emits a single
    # filler word ("photograph", "image") or a known-hallucinated
    # opener on a frame that's too dark / too noisy to parse. Those
    # captions pollute the rollup (Qwen treats them as facts). Reject
    # captions that are too short OR start with a common hallucination
    # prefix and aren't followed by substantive content. The good
    # captions on adjacent frames usually carry the same context.
    _HALLUCINATION_PREFIXES = (
        "a picture of ",
        "a photo of ",
        "an image of ",
        "screenshot of ",
        "the image shows ",
    )

    def _caption_is_useful(cap: str) -> bool:
        if not cap:
            return False
        stripped = cap.strip().rstrip(".")
        # Word-count gate. Empirically, captions <5 words on a
        # `<MORE_DETAILED_CAPTION>` task are filler.
        words = stripped.split()
        if len(words) < 5:
            return False
        # Reject hallucination prefixes ONLY when they're nearly the
        # whole caption (i.e. "a picture of a man" → reject; "a
        # picture of a man holding a clipboard" → keep).
        low = stripped.lower()
        for prefix in _HALLUCINATION_PREFIXES:
            if low.startswith(prefix) and len(words) < 8:
                return False
        return True

    captions: list[str] = []
    dropped_low_quality = 0
    for frame in frames:
        cap = _florence_caption(frame)
        if not cap:
            continue
        cleaned = cap.strip().rstrip(".")
        if _caption_is_useful(cleaned):
            captions.append(cleaned)
        else:
            dropped_low_quality += 1

    # Sprint I#6 — collapse near-duplicate captions before the Qwen
    # rollup ("a man in a suit" / "a man in a dark suit" → one). Keeps
    # the prompt tight and stops a repeated observation from skewing
    # the summary.
    captions_before_dedup = len(captions)
    captions = dedup_captions(captions)
    dropped_dupe_captions = captions_before_dedup - len(captions)

    # Audio transcription — captures what was SAID, which is usually
    # more informative than what the visuals show on talking-head /
    # interview / lecture clips. Failure modes (whisper not installed,
    # no audio track, transcription crash) return None and the rest
    # of the pipeline ignores it.
    # Sprint I#6 — does the container even HAVE an audio stream? Probed
    # up front so we can tell a genuinely-silent video (no audio track,
    # empty transcript expected) apart from a video that HAS audio but
    # whose transcription came back empty (a real failure worth
    # auditing). None when ffprobe is unavailable.
    has_audio_track = _video_has_audio_track(raw_bytes)

    transcript: str | None = None
    # Skip transcription entirely when we know there's no audio track —
    # saves a whisper invocation on silent screen-recordings / GIFs-as-
    # mp4 / b-roll, which are common in the dataset.
    if has_audio_track is False:
        transcript = None
    else:
        try:
            from backend.transcribe import transcribe_video_audio
            transcript = transcribe_video_audio(raw_bytes)
        except Exception:
            logger.exception("video: transcribe step crashed")
            transcript = None
    # Cap the transcript size before sending to Qwen. Whisper can
    # emit kilobytes of text on a long video; Qwen's context window
    # is finite and we'd rather spend it on the per-frame captions
    # too. ~1500 chars ≈ 250 words of speech, plenty to anchor the
    # subject of a video.
    if transcript and len(transcript) > 1500:
        transcript = transcript[:1500].rsplit(" ", 1)[0] + "…"

    # Filename → fallback topic. Strip the extension, swap _/- for
    # spaces, condense whitespace.
    fname_topic = ""
    if image.original_filename:
        fname_topic = (
            image.original_filename.rsplit(".", 1)[0]
            .replace("_", " ")
            .replace("-", " ")
        ).strip()
        fname_topic = re.sub(r"\s+", " ", fname_topic)

    # Pack visual captions + spoken transcript into one Qwen call.
    # `caption` carries the per-keyframe visual observations;
    # `ocr_text` field is repurposed to carry the audio transcript
    # because the Qwen prompt already treats it as authoritative
    # extracted text (the "USER EXPERIENCE AND AI TOOLS" on-screen
    # caption Florence picked up still lands here too, just merged
    # in with the spoken content). Qwen weighs both when synthesizing
    # — spoken content typically dominates when present because it's
    # longer and carries the actual subject of the video.
    summary: Optional[str] = None
    if captions or transcript:
        joined = " | ".join(captions) if captions else ""
        caption_for_llm = (
            f"Video — observed across {len(captions)} keyframes: {joined}"
            if captions
            else "Video — no visual frames captured."
        )
        ocr_for_llm = (
            f"Spoken content (transcribed): {transcript}"
            if transcript
            else None
        )
        try:
            summary = _llm_rewrite_summary(
                caption=caption_for_llm,
                names=[],
                ocr_text=ocr_for_llm,
                scene=None,
                setting=None,
                content_type="video",
                tags=None,
                regions=None,
                objects=None,
                concepts=None,
                vlm_description=None,
            )
        except Exception:
            logger.exception("video: qwen rewrite failed; using raw captions")
            summary = None
        if not summary:
            # Qwen disabled or failed — prefer the transcript when we
            # have one (most information per character), then the
            # longest single caption.
            if transcript:
                summary = transcript
            elif captions:
                summary = max(captions, key=len)
    if not summary:
        summary = "Video file. Preview unavailable."

    # Topic generation: derive from the summary's first noun phrase
    # when we can. The simplest heuristic that works well: take the
    # first sentence and Title-Case it, capped at ~6 words. Falls
    # back to the cleaned filename.
    topic = _derive_video_topic(summary, fname_topic)

    points: list[str] = []
    if duration_s and duration_s > 0:
        mins, secs = divmod(int(duration_s), 60)
        if mins > 0:
            points.append(f"Duration: {mins}m {secs}s")
        else:
            points.append(f"Duration: {secs}s")
    if image.byte_size_original:
        mb = image.byte_size_original / (1024 * 1024)
        points.append(f"Size: {mb:.1f} MB")
    if captions:
        points.append(f"Sampled {len(captions)} keyframes")

    cat = _classify_content(summary, image.original_filename, "video")

    # Record summary_signals on the image so quality regressions are
    # debuggable post-hoc. Same `image.__dict__["summary_signals"] =
    # ...` pattern the image pipeline uses (see ~line 567) — avoids
    # SQLAlchemy state-load on an expired column and the async-greenlet
    # crash that triggers. The caller's `_mark_done` reads this dict
    # and writes it as JSONB.
    signals: dict = {
        "kind": "video",
        "duration_s": float(duration_s) if duration_s else None,
        "frame_count": len(frames),
        # Sprint I#6 — scene-cut + caption-dedup telemetry, so a
        # quality regression ("summary missed the second half of the
        # video") is debuggable: did we under-sample scenes, or did
        # dedup collapse too aggressively?
        "frames_sampled": frames_before_scenecut,
        "scenes_kept": scenes_kept,
        "caption_count": len(captions),
        "dropped_low_quality_captions": dropped_low_quality,
        "dropped_dupe_captions": dropped_dupe_captions,
        # Sprint I#6 — audio-presence signal. Distinguishes silent
        # video from failed transcription (see _video_has_audio_track).
        "has_audio_track": has_audio_track,
        "has_transcript": bool(transcript),
        "transcript_chars": len(transcript) if transcript else 0,
        # True only when the video HAS audio but transcription produced
        # nothing — the row an operator should look at.
        "transcription_gap": bool(has_audio_track) and not bool(transcript),
        "qwen_succeeded": bool(summary) and summary not in (
            "Video file. Preview unavailable.",
        ),
    }
    # Merge with any signals the caller already set so we don't clobber.
    existing = image.__dict__.get("summary_signals") or {}
    if isinstance(existing, dict):
        existing.update(signals)
        image.__dict__["summary_signals"] = existing
    else:
        image.__dict__["summary_signals"] = signals

    return SummaryResult(
        topic=topic, summary=summary, points=points[:5], content_type=cat,
    )


def _derive_video_topic(summary: str, fallback: str) -> str:
    """Short, title-cased noun phrase for the gallery card.

    Trim the summary to its first sentence, drop opening filler
    ("A video showing", "The video depicts"), Title-Case, cap at
    6 words. Falls back to `fallback` (typically the cleaned
    filename) when the summary doesn't yield anything useful.
    """
    if not summary:
        return fallback or "Video"
    first = re.split(r"[.!?]", summary, maxsplit=1)[0].strip()
    if not first:
        return fallback or "Video"
    # Strip openers — Qwen often produces "A video showing X" or
    # "The clip depicts Y." The opener is conversational filler,
    # not information.
    first = re.sub(
        r"^(?:a |the )?(?:video |clip |scene |shot |sequence |footage )?"
        r"(?:that )?(?:is )?(?:showing|depicts|features|captures|shows)\s+",
        "",
        first,
        flags=re.IGNORECASE,
    )
    # Cap words.
    words = first.split()
    if len(words) > 6:
        words = words[:6]
    topic = " ".join(w[0].upper() + w[1:] if w else w for w in words)
    return topic or (fallback or "Video")


# --- audio -----------------------------------------------------------------


def _summarize_audio(image: Image, raw_bytes: bytes) -> Optional[SummaryResult]:
    """Standalone audio-file summary (mp3 / m4a / wav / flac / ogg / …).

    Audio rows have no video stream and no keyframes, so there's nothing
    for Florence to caption — the *entire* signal is what was SAID. The
    flow mirrors the transcript half of `_summarize_video`:

      1. Transcribe the audio via `transcribe_video_audio` (it runs
         `ffmpeg -vn -ac 1 -ar 16000` then faster-whisper — `-vn` is a
         no-op on an audio-only input, so it transcribes mp3/m4a/wav/
         flac/ogg the same way it does a video's audio track).
      2. Feed the transcript through the SAME Qwen rewrite the video path
         uses (`_llm_rewrite_summary`, content_type="audio") so the
         result is a concise CONTENT summary — the subject / speakers /
         topics — not a verbatim transcript dump.
      3. Derive a short title via `_derive_video_topic` from that summary.

    Every step degrades independently so an install without the [ml]
    extras (no whisper) or without ffmpeg never crashes the worker and
    never falls back to echoing the UUID filename:
      - whisper unavailable / ffmpeg missing → transcript is None
      - audio is silent / instrumental music → transcript is "" / None
      - Qwen disabled or failed            → use the raw transcript
      - no speech at all                   → a generic but sensible
        "Audio recording — no spoken content detected" (+ duration).
    """
    # Duration (ffprobe reads `format=duration` — works on audio
    # containers just as it does on video). Used for the fallback
    # summary + the `points` line. None when ffprobe is unavailable.
    duration_s = _probe_video_duration(raw_bytes)

    # Does the file actually carry a decodable audio stream? ffprobe's
    # audio-stream select works on mp3/m4a/wav/flac too. None when
    # ffprobe is unavailable — in that case we still attempt the
    # transcribe (whisper itself will fail gracefully). This lets us
    # tell a genuinely-silent / non-audio blob apart from a transcription
    # that came back empty despite real audio (worth auditing).
    has_audio_track = _video_has_audio_track(raw_bytes)

    transcript: str | None = None
    if has_audio_track is False:
        # No audio stream — don't bother spinning up whisper.
        transcript = None
    else:
        try:
            from backend.transcribe import transcribe_video_audio
            transcript = transcribe_video_audio(raw_bytes)
        except Exception:
            logger.exception("audio: transcribe step crashed")
            transcript = None

    # Cap the transcript before Qwen — same budget the video path uses
    # so a long podcast doesn't blow the context window. ~1500 chars ≈
    # 250 words, plenty to anchor the subject.
    if transcript and len(transcript) > 1500:
        transcript = transcript[:1500].rsplit(" ", 1)[0] + "…"

    # Filename → fallback topic (extension stripped, _/- → spaces). Only
    # ever used as the *topic* fallback, never as the summary body, so we
    # never echo the raw UUID filename as the description.
    fname_topic = ""
    if image.original_filename:
        fname_topic = (
            image.original_filename.rsplit(".", 1)[0]
            .replace("_", " ")
            .replace("-", " ")
        ).strip()
        fname_topic = re.sub(r"\s+", " ", fname_topic)

    # Human-readable duration for the fallback summary line.
    def _dur_phrase() -> str:
        if not (duration_s and duration_s > 0):
            return ""
        mins, secs = divmod(int(duration_s), 60)
        if mins > 0:
            return f" ({mins}m {secs}s)"
        return f" ({secs}s)"

    summary: Optional[str] = None
    if transcript:
        # Repurpose the `ocr_text` field to carry the spoken transcript,
        # exactly as `_summarize_video` does — the Qwen prompt already
        # treats that field as authoritative spoken content. There are no
        # visual captions for audio, so `caption` just states the medium.
        #
        # IMPORTANT: pass content_type="video" (NOT "audio") here. Inside
        # `_llm_rewrite_summary`, content_type only selects the PROMPT /
        # transcript-label / length-caps — and the "video" branch is the
        # one wired for spoken content: it labels ocr_text as "Spoken
        # content (what was said)", uses the paraphrase-don't-quote
        # instructions that turn a transcript into a CONTENT summary, and
        # raises the reply cap to 1300 chars / 220 tokens so a long
        # transcript isn't rejected back to a verbatim dump. The "audio"
        # value would fall through to the IMAGE branch (a 700-char,
        # describe-the-visuals prompt) and reject our long transcript,
        # leaving a raw transcript echo. This arg is NOT persisted — the
        # row's content_type comes from `_classify_content` below.
        try:
            summary = _llm_rewrite_summary(
                caption="Audio-only recording — no video; describe from the spoken content.",
                names=[],
                ocr_text=f"Spoken content (transcribed): {transcript}",
                scene=None,
                setting=None,
                content_type="video",
                tags=None,
                regions=None,
                objects=None,
                concepts=None,
                vlm_description=None,
            )
        except Exception:
            logger.exception("audio: qwen rewrite failed; using raw transcript")
            summary = None
        if not summary:
            # Qwen disabled or failed — the transcript itself is the most
            # information-dense thing we have. Still beats a filename echo.
            summary = transcript
        elif summary:
            # We reused the VIDEO prompt to get transcript-aware
            # summarization, so Qwen sometimes frames the output as
            # "In this video, the speaker discusses…" / "The video
            # covers…". This is an AUDIO file with no visuals, so rewrite
            # those video references to audio-appropriate wording. Done as
            # a post-step (not by forking the shared prompt) so real video
            # summaries are unaffected.
            summary = re.sub(
                r"\bin this video\b", "In this audio recording",
                summary, count=1, flags=re.IGNORECASE,
            )
            summary = re.sub(
                r"\b(?:this|the) video\b", "this recording",
                summary, flags=re.IGNORECASE,
            )
            summary = re.sub(r"\bvideo\b", "recording", summary, flags=re.IGNORECASE)
            summary = re.sub(r"\b(?:viewers|the viewer)\b", "listeners", summary, flags=re.IGNORECASE)
            # Re-capitalize the leading char in case a substitution landed
            # at position 0.
            if summary and summary[0].islower():
                summary = summary[0].upper() + summary[1:]

    if not summary:
        # No speech detected: silent file, instrumental music, whisper
        # unavailable, or extraction failed. Emit a sensible, honest
        # description (+ duration) rather than the UUID filename.
        summary = f"Audio recording — no spoken content detected{_dur_phrase()}."

    # Short title from the summary; falls back to the cleaned filename,
    # then a literal "Audio". Strip the conversational opener Qwen tends
    # to produce ("In this audio recording, the speaker discusses …")
    # FIRST so `_derive_video_topic` title-cases the actual subject rather
    # than the filler lead-in. Whatever's left after the opener is the
    # topic seed; if the regex doesn't match, the summary passes through
    # unchanged and `_derive_video_topic` handles the "A video showing X"
    # style openers it already knows.
    topic_seed = re.sub(
        r"^(?:in\s+)?(?:this|the)\s+(?:audio\s+)?recording[,:]?\s*"
        r"(?:the\s+)?(?:speaker|host|narrator|presenter|person)?\s*"
        r"(?:discusses|describes|explains|covers|talks about|"
        r"goes over|walks through|presents|outlines)?\s*",
        "",
        summary,
        count=1,
        flags=re.IGNORECASE,
    )
    topic = _derive_video_topic(topic_seed or summary, fname_topic or "Audio")

    points: list[str] = []
    if duration_s and duration_s > 0:
        mins, secs = divmod(int(duration_s), 60)
        if mins > 0:
            points.append(f"Duration: {mins}m {secs}s")
        else:
            points.append(f"Duration: {secs}s")
    if image.byte_size_original:
        mb = image.byte_size_original / (1024 * 1024)
        points.append(f"Size: {mb:.1f} MB")
    if transcript:
        points.append("Transcribed spoken content")

    # Classify with the video taxonomy — it carries the spoken-content
    # buckets that fit audio (music / presentation / lecture / vlog /
    # interview), whereas the doc taxonomy is paper-oriented.
    cat = _classify_content(summary, image.original_filename, "video")

    # Telemetry mirroring the video path's `summary_signals`, so a
    # quality regression on audio rows is debuggable post-hoc. Same
    # `image.__dict__[...]` write pattern (avoids the async-greenlet
    # column-load crash); `_mark_done` persists it as JSONB.
    signals: dict = {
        "kind": "audio",
        "duration_s": float(duration_s) if duration_s else None,
        "has_audio_track": has_audio_track,
        "has_transcript": bool(transcript),
        "transcript_chars": len(transcript) if transcript else 0,
        # True only when audio IS present but transcription produced
        # nothing — the row an operator should look at.
        "transcription_gap": bool(has_audio_track) and not bool(transcript),
        "qwen_succeeded": bool(transcript) and summary != transcript
        and not summary.startswith("Audio recording — no spoken content"),
    }
    existing = image.__dict__.get("summary_signals") or {}
    if isinstance(existing, dict):
        existing.update(signals)
        image.__dict__["summary_signals"] = existing
    else:
        image.__dict__["summary_signals"] = signals

    return SummaryResult(
        topic=topic, summary=summary, points=points[:5], content_type=cat,
    )


def _probe_video_duration(raw_bytes: bytes) -> Optional[float]:
    """Use ffprobe to read duration in seconds. Returns None on any
    failure — the caller substitutes fixed offsets in that case.

    Same temp-file approach as `_extract_keyframe`: piping through
    stdin breaks on mp4s whose moov atom lives at the END of the file
    (most phone recordings, screen captures). ffprobe with a real
    file path can seek freely to the moov atom regardless of where
    it sits in the container.
    """
    import tempfile
    import os
    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix="probe-", suffix=".bin", delete=False,
        ) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name
        from backend.ffmpeg_args import safe_input_args
        proc = subprocess.run(
            [
                "ffprobe",
                "-loglevel", "error",
                *safe_input_args(),
                "-show_entries", "format=duration",
                "-of", "csv=p=0",
                tmp_path,
            ],
            capture_output=True,
            timeout=15,
        )
        if proc.returncode != 0:
            return None
        s = proc.stdout.decode(errors="replace").strip()
        return float(s) if s else None
    except (FileNotFoundError, subprocess.TimeoutExpired, ValueError, OSError):
        return None
    finally:
        if tmp_path:
            try: os.unlink(tmp_path)
            except OSError: pass


def _extract_keyframe(raw_bytes: bytes, seek_seconds: float = 5) -> Optional[bytes]:
    """Write bytes to a temp file, run ffmpeg against the file path, capture one PNG.

    Requires `ffmpeg` on PATH. Returns None on any failure (missing binary,
    unsupported container, seek past end of file, etc.) so the caller can
    degrade gracefully.

    **Why a temp file instead of piping through stdin (2026-05 fix):**
    Most consumer MP4s (phone recordings, screen captures, content from
    Drive that wasn't run through `+faststart`) store the `moov` atom at
    the END of the file. When ffmpeg reads from a pipe, it can't seek
    backwards to that moov atom — the container parse fails with
    `partial file` at offset 0x30. The user's library had this exact
    issue on every video that went through HLS transcoding, because
    after HLS the `original_blob_key` still pointed at the original
    mp4 but the worker was piping it through stdin → ffmpeg failed →
    every video got the generic fallback summary "Video file. Preview
    unavailable." Writing to a temp file lets ffmpeg seek freely.

    Uses **post-input** `-ss` (input file, then seek). The previous
    pre-input form is faster but can land between keyframes on
    poorly-indexed MP4s and return EMPTY stdout with exit code 0.
    Post-input is decoder-accurate.
    """
    import tempfile
    import os
    tmp_path: Optional[str] = None
    try:
        # NamedTemporaryFile with delete=False so ffmpeg can open the
        # path on Windows + we control the cleanup in the finally.
        with tempfile.NamedTemporaryFile(
            prefix="kf-", suffix=".bin", delete=False,
        ) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name
        from backend.ffmpeg_args import safe_input_args
        proc = subprocess.run(
            [
                "ffmpeg",
                "-loglevel", "error",
                *safe_input_args(),
                "-i", tmp_path,
                "-ss", f"{seek_seconds:.3f}",
                "-frames:v", "1",
                "-f", "image2pipe",
                "-vcodec", "png",
                "pipe:1",
            ],
            capture_output=True,
            timeout=30,
            check=True,
        )
        if not proc.stdout:
            # exit 0 with empty payload: seek landed past the last
            # frame, or container has no decodable video stream.
            logger.info(
                "extract_keyframe: ffmpeg returned 0 bytes at seek=%.3fs; stderr=%r",
                seek_seconds, proc.stderr.decode(errors="replace")[:200],
            )
            return None
        return proc.stdout
    except subprocess.CalledProcessError as exc:
        logger.info(
            "extract_keyframe: ffmpeg exit %d at seek=%.3fs; stderr=%r",
            exc.returncode, seek_seconds,
            (exc.stderr or b"").decode(errors="replace")[:200],
        )
        return None
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return None
    finally:
        if tmp_path:
            try: os.unlink(tmp_path)
            except OSError: pass


# --- video summary Batch 2 (Sprint I#6) ------------------------------------


def _video_has_audio_track(raw_bytes: bytes) -> Optional[bool]:
    """ffprobe whether the video carries an audio stream.

    Returns True/False, or None if ffprobe is unavailable / errors.
    Lets the summary distinguish three cases that the old single
    `has_transcript` flag conflated:
      - has_audio_track == False           → genuinely silent video;
        an empty transcript is EXPECTED, not a failure.
      - has_audio_track == True,
        has_transcript == False            → transcription failed or
        produced nothing despite audio being present — worth auditing.
      - both True                          → healthy.
    """
    import tempfile
    import os
    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix="aprobe-", suffix=".bin", delete=False,
        ) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name
        from backend.ffmpeg_args import safe_input_args
        proc = subprocess.run(
            [
                "ffprobe",
                "-loglevel", "error",
                *safe_input_args(),
                "-select_streams", "a",
                "-show_entries", "stream=codec_type",
                "-of", "csv=p=0",
                tmp_path,
            ],
            capture_output=True,
            timeout=15,
        )
        if proc.returncode != 0:
            return None
        return b"audio" in proc.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        return None
    finally:
        if tmp_path:
            try: os.unlink(tmp_path)
            except OSError: pass


def _frame_histogram(png_bytes: bytes) -> Optional[list[float]]:
    """Coarse 256-bin luminance histogram of a frame, L1-normalized.
    Returns None if the frame can't be decoded."""
    try:
        from PIL import Image as _PILImage  # type: ignore
        img = _PILImage.open(BytesIO(png_bytes)).convert("L").resize((64, 64))
        hist = img.histogram()  # 256 bins for an "L" image
        total = float(sum(hist)) or 1.0
        return [c / total for c in hist]
    except Exception:
        return None


def _hist_distance(a: list[float], b: list[float]) -> float:
    """L1 distance between two normalized histograms — 0 (identical)
    to 2 (no overlap)."""
    return sum(abs(x - y) for x, y in zip(a, b))


def dedup_frames_by_histogram(
    frames: list[bytes], threshold: float = 0.35,
) -> list[bytes]:
    """Scene-cut detection by histogram diff (Sprint I#6).

    The video sampler extracts frames evenly across the duration, so a
    static talking-head clip yields ~30 near-identical frames — all
    captioned by Florence (expensive) and all fed to Qwen (redundant).
    This keeps a frame only when its luminance histogram differs from
    the LAST KEPT frame by >= `threshold`, biasing the kept set toward
    information-dense moments (actual scene changes) and cutting
    Florence calls on dead-static footage.

    Always keeps the first frame. Frames that can't be decoded are
    kept (can't compare → don't drop). Threshold 0.35 on the 0–2 L1
    scale empirically separates "same scene, camera shake" from "the
    shot changed."
    """
    if len(frames) <= 1:
        return frames
    kept: list[bytes] = []
    last_hist: Optional[list[float]] = None
    for f in frames:
        hist = _frame_histogram(f)
        if hist is None:
            kept.append(f)
            continue
        if last_hist is None or _hist_distance(hist, last_hist) >= threshold:
            kept.append(f)
            last_hist = hist
    return kept


def dedup_captions(captions: list[str], threshold: float = 0.82) -> list[str]:
    """Collapse near-duplicate keyframe captions before the Qwen rollup
    (Sprint I#6).

    Florence emits slight variations of the same observation across
    adjacent frames — "a man in a suit" / "a man in a dark suit" /
    "a man wearing a suit standing". Feeding all three to Qwen as
    distinct facts inflates the prompt and biases the summary toward
    whatever got repeated. We keep one representative (the longest,
    usually the most descriptive) per near-duplicate cluster, using a
    difflib similarity ratio.
    """
    import difflib
    kept: list[str] = []
    for cap in captions:
        low = cap.lower()
        matched = False
        for i, existing in enumerate(kept):
            ratio = difflib.SequenceMatcher(None, low, existing.lower()).ratio()
            if ratio >= threshold:
                # Keep the longer (more descriptive) of the pair.
                if len(cap) > len(existing):
                    kept[i] = cap
                matched = True
                break
        if not matched:
            kept.append(cap)
    return kept


# --- code ------------------------------------------------------------------
#
# Source files arrive as category="document" with a `text/x-<lang>` MIME
# (upload_validation._CODE_EXTS). Before #184 the doc path only routed a
# handful of structured-text extensions (.txt/.md/.csv/.json/...) through
# the plaintext extractor, so a `.py`/`.js`/`.go`/`.java` file fell
# through to the thin "Document: <stem>." stub — no language, no purpose,
# nothing searchable. We now detect the language from the extension /
# MIME and infer the file's PURPOSE from its imports + top-level defs +
# docstring, producing summaries like
#   "Python script for machine-learning model training (uses pytorch,
#    numpy, pandas)."
# that semantic search can surface by topic, not just by filename.

# Extension → human-readable language name. Mirrors the curated
# `_CODE_EXTS` map in upload_validation but maps to a display language
# rather than a MIME. Kept deliberately broad so the common languages
# all get a real "<Language> …" lead instead of the generic stub.
_CODE_EXT_LANG: dict[str, str] = {
    "py": "Python", "pyi": "Python", "ipynb": "Jupyter notebook",
    "js": "JavaScript", "mjs": "JavaScript", "cjs": "JavaScript",
    "jsx": "JavaScript (React)", "ts": "TypeScript", "tsx": "TypeScript (React)",
    "vue": "Vue", "svelte": "Svelte", "astro": "Astro",
    "rb": "Ruby", "php": "PHP", "java": "Java", "kt": "Kotlin",
    "kts": "Kotlin", "scala": "Scala", "swift": "Swift", "go": "Go",
    "rs": "Rust", "c": "C", "h": "C header", "cpp": "C++", "cc": "C++",
    "cxx": "C++", "hpp": "C++ header", "cs": "C#", "dart": "Dart",
    "lua": "Lua", "r": "R", "pl": "Perl", "sh": "Shell script",
    "bash": "Shell script", "zsh": "Shell script", "fish": "Shell script",
    "ps1": "PowerShell script", "sql": "SQL", "clj": "Clojure",
    "cljs": "ClojureScript", "ex": "Elixir", "exs": "Elixir",
    "elm": "Elm", "erl": "Erlang", "hs": "Haskell", "ml": "OCaml",
    "fs": "F#", "nim": "Nim", "zig": "Zig", "cr": "Crystal",
    "groovy": "Groovy", "gradle": "Gradle build", "jl": "Julia",
    "sol": "Solidity", "tf": "Terraform", "hcl": "HCL", "nix": "Nix",
    "pas": "Pascal", "f90": "Fortran", "asm": "Assembly", "s": "Assembly",
    "vala": "Vala", "hx": "Haxe", "rkt": "Racket", "scm": "Scheme",
    "lisp": "Lisp", "el": "Emacs Lisp", "coffee": "CoffeeScript",
    "m": "Objective-C / MATLAB", "mm": "Objective-C++",
    "css": "CSS", "scss": "SCSS stylesheet", "sass": "Sass stylesheet",
    "less": "Less stylesheet", "html": "HTML", "htm": "HTML",
    "xml": "XML", "yaml": "YAML config", "yml": "YAML config",
    "toml": "TOML config", "ini": "INI config", "cfg": "config",
    "conf": "config", "properties": "properties config", "env": "environment config",
    "dockerfile": "Dockerfile", "makefile": "Makefile", "cmake": "CMake build",
    "graphql": "GraphQL schema", "gql": "GraphQL schema", "proto": "Protocol Buffers schema",
    "diff": "diff / patch", "patch": "diff / patch", "tex": "LaTeX",
    "rst": "reStructuredText", "adoc": "AsciiDoc",
    "bat": "Batch script", "cmd": "Batch script", "vim": "Vimscript",
    "tcl": "Tcl", "awk": "AWK script", "lean": "Lean", "agda": "Agda",
}

# Extensions / language names we treat as CODE (purpose-inferred) rather
# than prose documents. Plain text / markdown / data files (.txt, .md,
# .csv, .json, .log) keep the existing prose-document LLM path because
# their bodies summarize well as natural language. Everything in
# `_CODE_EXT_LANG` that isn't pure config/markup data is code-like.
_NON_CODE_TEXT_EXTS: frozenset[str] = frozenset({
    "txt", "md", "markdown", "csv", "tsv", "json", "jsonc", "json5",
    "log", "rst", "adoc", "asciidoc", "tex", "latex",
})

# import-statement patterns per language family → pull the imported
# module / package names out of a code file's head so we can name the
# libraries it uses ("uses pytorch, numpy"). Best-effort regexes; a
# language we don't have a pattern for just yields no libs (the summary
# still leads with the language + any top-level defs).
_IMPORT_PATTERNS: tuple[tuple[str, str], ...] = (
    # Python:  import x / from x import y
    ("python", r"^\s*(?:from|import)\s+([a-zA-Z_][\w.]*)"),
    # JS/TS:   import ... from 'x'  /  require('x')
    ("js", r"""(?:from|require\()\s*['"]([^'"]+)['"]"""),
    # Go:      import "x"
    ("go", r"""^\s*(?:import\s+)?["']([\w./-]+)["']"""),
    # Java/Kotlin/Scala:  import a.b.c
    ("java", r"^\s*import\s+(?:static\s+)?([a-zA-Z_][\w.]*)"),
    # Rust:    use a::b
    ("rust", r"^\s*use\s+([a-zA-Z_][\w:]*)"),
    # C/C++:   #include <x> / #include "x"
    ("c", r"""^\s*#\s*include\s*[<"]([^>"]+)[>"]"""),
    # Ruby:    require 'x'
    ("ruby", r"""^\s*require(?:_relative)?\s+['"]([^'"]+)['"]"""),
)

# Library / framework name → the kind of work it signals. Lets the
# heuristic say "for machine-learning model training" / "web server" /
# "data analysis" when the LLM is unavailable, instead of only listing
# raw module names. Substring match against the lower-cased import set.
_LIB_PURPOSE: tuple[tuple[tuple[str, ...], str], ...] = (
    (("torch", "tensorflow", "keras", "sklearn", "scikit", "xgboost",
      "lightgbm", "transformers", "jax", "fastai", "onnx", "openvino"),
     "machine-learning / AI model work"),
    (("numpy", "pandas", "scipy", "matplotlib", "seaborn", "polars",
      "statsmodels", "plotly"),
     "data analysis / scientific computing"),
    (("flask", "fastapi", "django", "express", "koa", "nestjs", "gin",
      "fiber", "actix", "rails", "sinatra", "spring", "starlette",
      "aiohttp", "tornado"),
     "a web server / API backend"),
    (("react", "vue", "angular", "svelte", "next", "nuxt", "preact",
      "solid-js"),
     "a web front-end / UI"),
    (("pytest", "unittest", "jest", "mocha", "vitest", "junit",
      "rspec", "cypress", "playwright", "selenium"),
     "automated tests"),
    (("requests", "httpx", "axios", "urllib", "aiohttp", "fetch"),
     "HTTP requests / API calls"),
    (("sqlalchemy", "psycopg2", "asyncpg", "pymongo", "redis",
      "sqlite3", "prisma", "mongoose", "gorm", "diesel"),
     "database access"),
    (("boto3", "google.cloud", "azure", "kubernetes", "docker",
      "terraform", "ansible"),
     "cloud / infrastructure automation"),
    (("argparse", "click", "typer", "cobra", "commander"),
     "a command-line tool"),
    (("selenium", "scrapy", "beautifulsoup", "bs4", "playwright",
      "puppeteer"),
     "web scraping / browser automation"),
    (("discord", "telegram", "slack", "telethon", "tweepy"),
     "a chat / social bot"),
    (("pygame", "arcade", "pyglet", "godot", "unity"),
     "a game"),
)

# Top-level definition patterns per language family — gives the summary a
# few concrete symbol names ("defines train_model, evaluate, DataLoader")
# which are exactly the terms a developer searches by.
_DEF_PATTERNS: tuple[tuple[str, str], ...] = (
    ("python", r"^\s*(?:async\s+)?def\s+([a-zA-Z_]\w*)"),
    ("python", r"^\s*class\s+([a-zA-Z_]\w*)"),
    ("js", r"^\s*(?:export\s+)?(?:default\s+)?(?:async\s+)?function\s+([a-zA-Z_$]\w*)"),
    ("js", r"^\s*(?:export\s+)?(?:abstract\s+)?class\s+([a-zA-Z_$]\w*)"),
    ("js", r"^\s*(?:export\s+)?const\s+([a-zA-Z_$]\w*)\s*=\s*(?:async\s*)?\("),
    ("go", r"^\s*func\s+(?:\([^)]*\)\s*)?([A-Za-z_]\w*)"),
    ("go", r"^\s*type\s+([A-Za-z_]\w*)\s+(?:struct|interface)"),
    ("rust", r"^\s*(?:pub\s+)?(?:async\s+)?fn\s+([a-zA-Z_]\w*)"),
    ("rust", r"^\s*(?:pub\s+)?(?:struct|enum|trait)\s+([a-zA-Z_]\w*)"),
    ("java", r"^\s*(?:public|private|protected)?\s*(?:static\s+)?(?:final\s+)?(?:class|interface|enum)\s+([A-Za-z_]\w*)"),
    ("ruby", r"^\s*(?:def|class|module)\s+([A-Za-z_]\w*)"),
)


def _code_language(filename: str, mime: str | None) -> Optional[str]:
    """Human language name for a code file, or None when it isn't code.

    Resolves by extension first (most reliable), then basename
    (Dockerfile / Makefile), then the `text/x-<lang>` MIME the upload
    pipeline assigned. Returns None for prose/data text the caller
    should keep on the normal document path.
    """
    name = (filename or "").lower()
    base = name.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
    ext = base.rsplit(".", 1)[-1] if "." in base else ""
    if ext in _NON_CODE_TEXT_EXTS:
        return None
    if ext and ext in _CODE_EXT_LANG:
        return _CODE_EXT_LANG[ext]
    # Basename match (no extension): Dockerfile, Makefile, …
    stem = base.split(".", 1)[0]
    if stem in _CODE_EXT_LANG:
        return _CODE_EXT_LANG[stem]
    # MIME fallback: text/x-python → "python", for a code file whose
    # extension we didn't recognize but whose upload MIME names the
    # language. Deliberately RESTRICTED to a known code-language allow-
    # list — a blind `text/x-*` → title-case would mis-route data text
    # MIMEs like `text/x-vcard` / `text/x-ics` / `text/calendar` to the
    # code path. Those return None here and fall through to the prose-
    # document summarizer instead.
    m = (mime or "").lower().split(";")[0].strip()
    _MIME_LANG = {
        "javascript": "JavaScript", "css": "CSS", "html": "HTML",
        "xml": "XML", "python": "Python", "typescript": "TypeScript",
        "shellscript": "Shell script", "csharp": "C#", "c++": "C++",
        "ruby": "Ruby", "php": "PHP", "java": "Java", "kotlin": "Kotlin",
        "scala": "Scala", "swift": "Swift", "go": "Go", "rust": "Rust",
        "c": "C", "sql": "SQL", "perl": "Perl", "lua": "Lua",
        "r": "R", "dart": "Dart", "haskell": "Haskell", "elixir": "Elixir",
        "clojure": "Clojure", "scss": "SCSS stylesheet", "sass": "Sass stylesheet",
        "less": "Less stylesheet", "powershell": "PowerShell script",
        "dockerfile": "Dockerfile", "makefile": "Makefile",
        "objectivec": "Objective-C", "toml": "TOML config",
        "yaml": "YAML config",
    }
    if m.startswith("text/x-") or m in (
        "text/javascript", "text/css", "text/html", "text/xml",
    ):
        token = m.split("/", 1)[-1].replace("x-", "")
        return _MIME_LANG.get(token)
    return None


def _family_for_lang(lang: str) -> str:
    """Map a display language to the regex-family key used by the
    import/def pattern tables ("Python (React)" → "python")."""
    low = lang.lower()
    if low.startswith("python") or "jupyter" in low:
        return "python"
    if any(k in low for k in ("javascript", "typescript", "vue", "svelte", "react")):
        return "js"
    if low.startswith("go"):
        return "go"
    if low.startswith("rust"):
        return "rust"
    if any(k in low for k in ("java", "kotlin", "scala")):
        return "java"
    if low == "c" or any(k in low for k in ("c++", "c header", "objective")):
        return "c"
    if low.startswith("ruby"):
        return "ruby"
    return ""


def _extract_code_facts(
    text: str, family: str,
) -> tuple[list[str], list[str], Optional[str]]:
    """Return (libraries, top_level_defs, leading_docstring/comment).

    All best-effort regex scans over the file head. `family` selects the
    import + def patterns; an unknown family yields empty lib/def lists
    (the summary still leads with the language). The docstring/comment is
    the first block comment or module docstring — often a one-line
    description of what the file does.
    """
    head = text[:8000]
    lines = head.splitlines()

    # Libraries — run only the patterns for this family (plus the C
    # include pattern which is unambiguous) over the head.
    libs: list[str] = []
    seen_lib: set[str] = set()
    for fam, pat in _IMPORT_PATTERNS:
        if fam != family:
            continue
        rx = re.compile(pat, re.MULTILINE)
        for mobj in rx.finditer(head):
            raw = (mobj.group(1) or "").strip()
            if not raw:
                continue
            if family == "go":
                # Go module paths: the LAST path segment is the package
                # name a developer recognizes (github.com/gin-gonic/gin →
                # gin; net/http → http; ./util → util).
                token = raw.strip("'\"").rstrip("/").split("/")[-1]
            elif family == "js":
                # JS specifiers: drop relative imports entirely (./foo,
                # ../bar — not libraries). For a scoped package keep the
                # scope's package (@tanstack/react-query → react-query);
                # otherwise the first path segment is the package
                # (lodash/debounce → lodash; react → react).
                spec = raw.strip("'\"")
                if spec.startswith("."):
                    continue
                if spec.startswith("@"):
                    parts = spec.split("/")
                    token = parts[1] if len(parts) > 1 else parts[0].lstrip("@")
                else:
                    token = spec.split("/")[0]
            else:
                # Python/Java/Rust/C: the FIRST segment is the top-level
                # package (numpy.linalg → numpy; a::b → a; stdio.h → stdio).
                token = raw.split(".")[0].split("/")[0].split("::")[0]
            token = token.strip().strip("'\"")
            low = token.lower()
            if not low or len(low) < 2 or low in seen_lib:
                continue
            # Skip relative / stdlib-noise imports that aren't useful tags.
            if low in {
                # Python stdlib noise
                "os", "sys", "re", "io", "abc", "std", "self",
                "__future__", "typing", "dataclasses", "enum",
                "functools", "itertools", "collections", "pathlib",
                "json", "math", "time", "datetime", "logging",
                "subprocess", "uuid", "random", "string", "warnings",
                # Go stdlib noise
                "fmt", "errors", "context", "strings", "strconv",
                "bytes", "bufio", "http", "net", "sort", "sync",
                # JS/Node stdlib / relative noise
                "path", "fs", "url", "util", "events", "stream",
            }:
                continue
            seen_lib.add(low)
            libs.append(token)
            if len(libs) >= 12:
                break

    # Top-level definitions.
    defs: list[str] = []
    seen_def: set[str] = set()
    for fam, pat in _DEF_PATTERNS:
        if fam != family:
            continue
        rx = re.compile(pat, re.MULTILINE)
        for mobj in rx.finditer(head):
            name = (mobj.group(1) or "").strip()
            low = name.lower()
            if not name or low in seen_def or len(name) < 2:
                continue
            seen_def.add(low)
            defs.append(name)
            if len(defs) >= 12:
                break

    # Leading docstring / comment — first meaningful prose in the file.
    doc: Optional[str] = None
    # Python / shell-style module docstring or # comment block.
    mdoc = re.search(r'^\s*(?:"""|\'\'\')(.+?)(?:"""|\'\'\')', head, re.DOTALL)
    if mdoc:
        doc = mdoc.group(1).strip()
    if not doc:
        # `#` is a COMMENT in Python/shell/Ruby/YAML but a PREPROCESSOR
        # directive in C/C++/Objective-C (`#include`, `#define`). Only
        # treat `#` as a comment marker for the non-C families so we
        # don't slurp `#include <stdio.h>` into the description.
        hash_is_comment = family != "c"
        markers = ["//", "/*", "*", ";;", "--"]
        if hash_is_comment:
            markers.append("#")
        comment_lines: list[str] = []
        for ln in lines[:40]:
            s = ln.strip()
            if s.startswith(tuple(markers)):
                cleaned = s.lstrip("#/*;-").strip()
                if cleaned and not cleaned.lower().startswith(
                    ("!", "-*-", "noqa", "type:", "pylint", "eslint",
                     "prettier", "coding:", "copyright", "spdx",
                     "include", "define", "pragma", "ifndef", "ifdef",
                     "endif", "import")
                ):
                    comment_lines.append(cleaned)
            elif s and comment_lines:
                break
            if len(comment_lines) >= 4:
                break
        if comment_lines:
            doc = " ".join(comment_lines)
    if doc:
        doc = re.sub(r"\s+", " ", doc).strip()
        if len(doc) > 300:
            doc = doc[:300].rsplit(" ", 1)[0] + "…"
    return libs, defs, doc


def _heuristic_code_summary(
    lang: str, libs: list[str], defs: list[str], doc: Optional[str],
    fname: str,
) -> str:
    """Build a useful, searchable code summary without the LLM.

    Composes "<Language> <kind> <purpose> (uses <libs>). Defines <defs>."
    so a developer can find the file by its job and its libraries even
    when Qwen is unavailable. Falls back to language + filename when we
    have nothing else — still better than "Document: <stem>."
    """
    # Purpose phrase from libraries, if any library family matches.
    low_libs = {l.lower() for l in libs}
    purpose = ""
    for keys, phrase in _LIB_PURPOSE:
        if any(any(k in lib for lib in low_libs) for k in keys):
            purpose = phrase
            break

    # "script" vs "module" vs "stylesheet" — small nicety from the
    # language label itself. Labels that already carry their own kind
    # (e.g. "Shell script", "YAML config", "Dockerfile") get no suffix.
    low = lang.lower()
    if low == "css":
        kind = "stylesheet"
    elif any(k in low for k in (
        "stylesheet", "config", "dockerfile", "makefile", "schema",
        "build", "script", "notebook", "header",
    )):
        # Label already carries its own kind word — don't double it.
        kind = ""
    else:
        kind = "source file"

    lead = lang if not kind else f"{lang} {kind}"
    sentence = lead
    if purpose:
        sentence += f" for {purpose}"
    if libs:
        sentence += f" (uses {', '.join(libs[:6])})"
    sentence = sentence.rstrip() + "."

    extra = ""
    if defs:
        extra = " Defines " + ", ".join(defs[:6]) + "."
    desc = ""
    if doc:
        # The file's own description leads when present — most reliable
        # statement of purpose. Keep it to one sentence.
        first = re.split(r"(?<=[.!?])\s", doc, maxsplit=1)[0].strip()
        if first and len(first) >= 12:
            desc = " " + (first if first[-1] in ".!?" else first + ".")

    out = (sentence + extra + desc).strip()
    if not out or out == ".":
        stem = _filename_stem(fname)
        return f"{lang} file{(' — ' + stem) if stem else ''}."
    return out


def _summarize_code(
    image: Image, text: str, lang: str,
) -> SummaryResult:
    """Summarize a source-code file: detect language + infer purpose.

    Prefers a code-aware Qwen rewrite (rich, natural) and falls back to
    the deterministic `_heuristic_code_summary` so the result is always
    useful and searchable. Tags = language + libraries so search can
    facet by stack ("python", "pytorch"); content_type stays "code".
    """
    fname = image.original_filename or ""
    family = _family_for_lang(lang)
    libs, defs, doc = _extract_code_facts(text, family)

    heuristic = _heuristic_code_summary(lang, libs, defs, doc, fname)

    # Try the LLM for a more natural one-liner. We feed it the language,
    # the extracted facts, and a head of the source — same Qwen the doc
    # path uses. On any failure we keep the heuristic (already useful).
    summary = _llm_code_summary(text, fname, lang, libs, defs, doc) or heuristic

    # Topic — short, search-friendly. Prefer an LLM topic; fall back to
    # "<Language> · <stem>".
    stem = _filename_stem(fname)
    topic = _llm_compose_doc_topic(text[:2000], fname) or (
        f"{lang} · {stem}" if stem else lang
    )

    # Tags: language token(s) + libraries. Lower-cased, deduped. These
    # are written by `_mark_done`'s adjective-tag path normally, but code
    # summaries rarely contain adjectives, so we surface the stack here
    # via summary_signals (search reads signals.concepts) AND ensure the
    # libs appear in the summary text (they already do) so FTS hits them.
    tag_set: list[str] = []
    seen: set[str] = set()
    for t in [lang.split()[0]] + libs:
        tl = t.lower().strip()
        if tl and tl not in seen:
            seen.add(tl)
            tag_set.append(tl)

    points: list[str] = []
    if defs:
        points.append("Defines: " + ", ".join(defs[:6]))
    if libs:
        points.append("Libraries: " + ", ".join(libs[:6]))
    loc = text.count("\n") + 1
    points.append(f"~{loc} lines")

    # Stash the code facts in summary_signals so search's haystack
    # (which reads signals.concepts) indexes the library + symbol names
    # even when they don't survive into the trimmed summary text.
    concepts = tag_set + [d.lower() for d in defs[:8]]
    image.__dict__["summary_signals"] = {
        "kind": "code",
        "language": lang,
        "libraries": libs[:12],
        "defs": defs[:12],
        "concepts": concepts[:20],
    }

    return SummaryResult(
        topic=topic or stem or lang,
        summary=summary,
        points=points[:5],
        content_type="code",
    )


def _llm_code_summary(
    text: str, filename: str, lang: str,
    libs: list[str], defs: list[str], doc: Optional[str],
) -> Optional[str]:
    """Qwen2.5-Instruct one-liner describing what a code file is FOR.

    Returns None when the rewriter is disabled / unavailable; the caller
    keeps the deterministic heuristic in that case. Guides the model with
    the pre-extracted facts so even a small checkpoint produces a
    purpose-first sentence ("Python script that trains a CNN image
    classifier with PyTorch") rather than echoing code lines.
    """
    if not settings.rewriter_enabled:
        return None
    try:
        import torch  # type: ignore

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        head = text[:6000]
        facts: list[str] = [f"Language: {lang}"]
        if libs:
            facts.append("Imports/libraries: " + ", ".join(libs[:12]))
        if defs:
            facts.append("Top-level definitions: " + ", ".join(defs[:12]))
        if doc:
            facts.append("File's own description/comment: " + doc)

        instructions = (
            "Write ONE short, search-friendly sentence (under 30 words) "
            "describing what this source-code file IS and what it is FOR "
            "— the kind of thing a developer would type into a search bar "
            "months later to find it. Start with the programming language "
            "and the file's role (script, module, component, test, config). "
            "Name the concrete purpose and the key libraries/frameworks it "
            "uses. Example: 'Python script for training a machine-learning "
            "image classifier using PyTorch and torchvision.' Do NOT quote "
            "code verbatim, do NOT list every function, and do NOT begin "
            "with 'This file' or 'The code'. Output only the sentence."
        )
        messages = [
            {
                "role": "system",
                "content": (
                    "You write concise, factual one-line descriptions of "
                    "source-code files. You never invent libraries or "
                    "behavior that isn't in the input."
                ),
            },
            {
                "role": "user",
                "content": (
                    instructions
                    + f"\n\nFilename: {filename}\n"
                    + "\n".join(facts)
                    + f"\n\nSource (head):\n{head}"
                ),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3500
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]
        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=80,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        reply = re.sub(r"\s+", " ", reply).strip().strip('"').strip("'").strip()
        reply = reply.lstrip("-•*").strip()
        # Reject body-text leaks (the model copied a code line) and
        # over-long rambles.
        if not reply or len(reply) > 400:
            return None
        if _looks_like_body_excerpt(reply, head):
            return None
        if reply[-1] not in ".!?":
            reply = reply.rstrip(",;:") + "."
        return reply[0].upper() + reply[1:]
    except Exception:
        logger.exception("qwen: code summary failed")
        return None


# --- document --------------------------------------------------------------


def _summarize_document(
    image: Image, raw_bytes: bytes
) -> Optional[SummaryResult]:
    fname = image.original_filename or ""

    # #184 — CODE files (`.py`/`.js`/`.go`/...) arrive as documents but
    # need language + purpose inference, not prose summarization. Detect
    # the language up front; when it IS code, extract its text and route
    # to the code summarizer. Non-code text / PDFs / Office docs fall
    # through to the existing prose-document path below.
    code_lang = _code_language(fname, image.mime_type_original)
    if code_lang is not None:
        try:
            code_text = _extract_code_source(fname, raw_bytes)
        except Exception:
            logger.exception("code extract failed for %s", fname)
            code_text = ""
        if code_text.strip():
            return _summarize_code(image, code_text, code_lang)
        # Empty / unreadable code file — still emit a language-led stub
        # so it's better than "Document: <stem>." and stays searchable.
        return SummaryResult(
            topic=(f"{code_lang} · {_filename_stem(fname)}"
                   if _filename_stem(fname) else code_lang),
            summary=_heuristic_code_summary(code_lang, [], [], None, fname),
            points=[],
            content_type="code",
        )

    text, topic = _extract_doc_text(fname, raw_bytes, image.mime_type_original)
    if not text.strip():
        return SummaryResult(
            topic=topic or _filename_stem(fname) or "Document",
            summary=_doc_summary_fallback(fname, topic),
            points=[],
        )

    # Sprint I D2 — stash per-chunk text + embeddings on the image
    # object so the async wrapper can persist them to document_chunks
    # after _mark_done. Embedding happens inside this sync thread to
    # reuse the already-loaded CLIP runtime; we don't pay an extra
    # model load. Best-effort — failures (CLIP missing, OOM) leave
    # the chunks list empty and the async wrapper skips persistence.
    try:
        chunk_texts = split_doc_for_embedding(text)
        if chunk_texts:
            from backend.vision.text_embed import embed_texts
            embeddings = embed_texts(chunk_texts)
            image.__dict__["doc_chunks"] = [
                {"chunk_index": i, "text": t, "embedding": e}
                for i, (t, e) in enumerate(zip(chunk_texts, embeddings))
            ]
    except Exception:
        logger.exception("doc chunks: split+embed failed for %s", fname)

    # We still cap at `summarize_doc_max_chars` (default 20 000) so a 200-
    # page PDF doesn't blow up generation memory, but the LLM path below
    # internally chunk-summarizes within that budget instead of skimming
    # the head.
    truncated = text[: settings.summarize_doc_max_chars]

    # Try the instruction-LLM first. Qwen2.5-Instruct (already in memory
    # for image rewriting) produces a short "what is this doc ABOUT"
    # blurb the user can search by from memory, instead of dumping the
    # first 400 chars of body text (the old fallback's failure mode).
    summary = _llm_doc_summary(truncated, fname, topic)
    if not summary:
        summary = _extractive_summary(truncated)

    # Reject body-text leaks. Both the small LLM and the sumy
    # extractive fallback occasionally produce output that's a verbatim
    # chunk of the input ("Part 1 — Inventory Think about or list…"
    # was the first paragraph of LBLF.pdf, not a description of it).
    # A search-friendly summary never has 60 consecutive characters
    # identical to the source. When it does, prefer the filename-stub
    # fallback so the preview panel doesn't quote the doc back to the
    # user as its own description.
    if summary and _looks_like_body_excerpt(summary, truncated):
        logger.info("doc-summary: body-text leak detected for %s, falling back", fname)
        summary = None

    # Don't ever emit raw body content as the description. If both the
    # LLM and the extractive paths failed, hand back a stub the user can
    # at least recognize. A wall of `truncated[:400]` is what made
    # LBLF.pdf's description quote the first questionnaire prompt in
    # full instead of summarizing the doc.
    summary = _clamp_doc_summary(summary) or _doc_summary_fallback(fname, topic)

    # Auto-derive points when the LLM ran (it produces a single paragraph,
    # not a bulleted list). Heuristic keypoints still work for docs that
    # already use markdown list syntax.
    points = _extract_keypoints(truncated)
    if not points and summary:
        points = _llm_keypoints(truncated, fname)

    # Always prefer the LLM-generated topic for documents — the
    # extraction heuristic (first 4-120-char line) produces things
    # like "LBLF" or page numbers. Falls back to extracted topic,
    # then filename stem.
    llm_topic = _llm_compose_doc_topic(truncated, fname)
    if llm_topic:
        topic = llm_topic
    elif not topic or len(topic) < 4 or topic.isdigit():
        topic = _llm_doc_topic(truncated, fname) or topic

    return SummaryResult(
        topic=topic or _filename_stem(fname) or "Document",
        summary=summary,
        points=points[:5],
        content_type=_classify_content(summary, fname, "document"),
    )


def _filename_stem(filename: str) -> str:
    if not filename:
        return ""
    stem = filename.rsplit(".", 1)[0]
    return stem.replace("_", " ").replace("-", " ").strip()


def _doc_summary_fallback(filename: str, topic: Optional[str]) -> str:
    """When we have no usable LLM output, write a recognizable stub from
    the filename + topic instead of quoting body text verbatim. Keeps
    the description short and searchable."""
    stem = _filename_stem(filename)
    if topic and 3 <= len(topic) <= 80:
        return f"Document about {topic}."
    if stem:
        return f"Document: {stem}."
    return "Document content could not be summarized."


def _clamp_doc_summary(summary: Optional[str]) -> Optional[str]:
    """Cap description length and strip leading body-text fingerprints.

    The doc preview panel only has room for ~3 lines. Long extractive
    output (sumy LSA pulling 4 sentences out of a 200-page contract)
    overflows the panel and is the opposite of "searchable by memory."
    Trim to ~280 chars at a sentence boundary so the description stays
    glanceable.
    """
    if not summary:
        return summary
    s = re.sub(r"\s+", " ", summary).strip()
    if not s:
        return None
    # Strip markdown / label prefixes Qwen sometimes prepends despite
    # the prompt asking for plain text:
    #   "Document Description:**"  /  "**Description:**"  /  "## Summary"
    # All of those leak the LLM's internal scaffolding into the
    # preview panel — strip leading labels and any orphan stars/hashes.
    s = re.sub(
        r"^(?:#+\s*|\*+\s*)?(?:document\s+)?(?:summary|description|overview|tl;dr)\s*[:\-—]\s*\*+\s*",
        "",
        s,
        flags=re.IGNORECASE,
    )
    s = s.lstrip("*# \t-—:").strip()
    if not s:
        return None
    if len(s) <= 280:
        return s
    head = s[:280]
    cut = max(head.rfind(". "), head.rfind("! "), head.rfind("? "))
    if cut > 120:
        return head[: cut + 1].strip()
    return head.rstrip(",;: ") + "…"


def _looks_like_body_excerpt(summary: str, body: str) -> bool:
    """True when the LLM regurgitated a chunk of the input verbatim.

    Qwen2.5-1.5B sometimes copies the opening lines of a document instead
    of paraphrasing them ("Part 1 — Inventory Think about or list…"
    instead of "Class assignment about an inventory of skills."). The
    extractive sumy fallback path also produces these. Detect by
    sliding-window: if any 60-char run from the summary appears
    verbatim in the body, treat it as a leak and reject.

    60 chars is wide enough to skip incidental phrase overlap ("the
    document", "first page") but narrow enough to catch any real
    body quote — a search-friendly description never has 60
    consecutive characters identical to the source.
    """
    if not summary or not body:
        return False
    s = re.sub(r"\s+", " ", summary).lower()
    b = re.sub(r"\s+", " ", body).lower()
    if len(s) < 60:
        return False
    window = 60
    for i in range(0, len(s) - window + 1, 20):
        if s[i : i + window] in b:
            return True
    return False


def _extract_doc_text(
    filename: str, raw: bytes, mime: str | None = None,
) -> tuple[str, Optional[str]]:
    """Return (text, topic). Topic falls back to None when nothing usable.

    #184 — broadened so ANY textual document (not just a hardcoded list)
    gets summarized: after the binary-format handlers (pdf/docx/xlsx) we
    fall back to plaintext extraction for anything whose MIME is `text/*`
    / a structured-text type, OR that decodes cleanly as UTF-8 text. This
    is what lets `.vcf` (vCard), `.ics` (calendar), `.rtf`, and other
    text formats produce a real summary instead of the "Document: <stem>."
    stub. Genuinely-binary blobs (a 3D `.stl`, a `.obj` mesh) still return
    ("", None) and the caller emits the filename fallback.
    """
    name = filename.lower()
    m = (mime or "").lower().split(";")[0].strip()
    try:
        if name.endswith(".pdf") or m == "application/pdf":
            return _extract_pdf(raw)
        if name.endswith(".docx") or m == (
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.document"
        ):
            return _extract_docx(raw)
        if name.endswith((".xlsx", ".xls")) or m in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            return _extract_xlsx(raw)
        if name.endswith((".pptx", ".ppt")) or m in (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ):
            return _extract_pptx(raw)
        # Known structured-text extensions (fast path — no decode probe).
        if name.endswith((".txt", ".md", ".csv", ".log", ".json", ".tsv",
                          ".yaml", ".yml", ".toml", ".ini", ".sql", ".html",
                          ".htm", ".xml", ".rtf", ".vcf", ".ics", ".rst")):
            return _extract_plaintext(raw, filename)
        # MIME-driven text fallback: any text/* (incl. text/x-vcard,
        # text/calendar, text/x-<anything>) is plaintext-extractable.
        if m.startswith("text/") or m in (
            "application/json", "application/xml", "application/x-yaml",
            "application/rtf",
        ):
            return _extract_plaintext(raw, filename)
        # Last resort: probe whether the bytes are UTF-8 text. If a high
        # fraction decode without replacement, treat as plaintext — this
        # catches text formats we have no extension/MIME mapping for at
        # all. Binary meshes / blobs fail the probe and fall through.
        if _looks_like_text(raw):
            return _extract_plaintext(raw, filename)
    except Exception:
        logger.exception("doc extract failed for %s", filename)
    return "", None


def _looks_like_text(raw: bytes, sample: int = 4096) -> bool:
    """Heuristic: do the first `sample` bytes look like UTF-8 text?

    A NUL byte is a strong binary signal (text files almost never contain
    one). Otherwise we decode the sample as UTF-8 and require that very
    few characters are replacement/control chars. Conservative on purpose
    — a false positive just produces a slightly noisy plaintext summary,
    while a false negative falls back to the harmless filename stub.
    """
    if not raw:
        return False
    chunk = raw[:sample]
    if b"\x00" in chunk:
        return False
    try:
        decoded = chunk.decode("utf-8")
    except UnicodeDecodeError:
        return False
    if not decoded:
        return False
    # Count control chars (excluding common whitespace). A real text file
    # has almost none.
    ctrl = sum(
        1 for c in decoded
        if ord(c) < 32 and c not in "\t\n\r\f\v"
    )
    return (ctrl / len(decoded)) < 0.02


def _extract_pdf(raw: bytes) -> tuple[str, Optional[str]]:
    """PDF text extraction. pypdf is the default; if it produces almost
    nothing (typical for layout-heavy / two-column / image-PDFs) we fall
    back to pdfminer.six which handles those better. Image-only PDFs
    return empty text — caller emits a "could not extract" summary."""
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(BytesIO(raw))
    pages_text = []
    for page in reader.pages[:50]:  # cap at 50 pages (was 30)
        try:
            pages_text.append(page.extract_text() or "")
        except Exception:
            continue
    text = "\n\n".join(pages_text).strip()

    # If pypdf bailed (returned <80 visible chars from a multi-page doc),
    # try pdfminer.six. It's slower but parses encoded fonts and column
    # layouts pypdf trips on.
    if len(text) < 80 and len(reader.pages) > 0:
        miner_text = _pdfminer_extract(raw)
        if miner_text and len(miner_text) > len(text):
            text = miner_text.strip()

    # If BOTH pypdf and pdfminer returned almost nothing, the PDF is
    # likely image-only (scanned docs, contracts photographed page-by-
    # page). Rasterize each page with PyMuPDF and run Florence-2 <OCR>
    # over the rasters. Slow (~3-5 s/page on GPU) but unlocks summaries
    # for previously-opaque PDFs.
    if len(text) < 80 and len(reader.pages) > 0:
        ocr_text = _pdf_ocr_fallback(raw)
        if ocr_text and len(ocr_text) > len(text):
            text = ocr_text.strip()

    topic: Optional[str] = None
    md = reader.metadata or {}
    if md and md.title:
        topic = str(md.title).strip() or None
    if not topic:
        # First non-blank line truncated.
        for line in text.splitlines():
            line = line.strip()
            if 4 <= len(line) <= 120:
                topic = line
                break
    return text, topic


def _pdfminer_extract(raw: bytes) -> str:
    """pdfminer.six fallback. Returns "" on any failure (incl. missing
    optional dep) so the caller keeps whatever pypdf produced."""
    try:
        from pdfminer.high_level import extract_text  # type: ignore

        return extract_text(BytesIO(raw)) or ""
    except Exception:
        return ""


def _pdf_pages_to_images(raw: bytes, max_pages: int = 10) -> list[bytes]:
    """Rasterize the first N pages of a PDF to PNG bytes via PyMuPDF.

    Returns a list of PNG bytes (empty list on failure). Cap of 10 pages
    is a cost guard — for the summary use case we don't need the whole
    doc, just enough text to produce a meaningful description. The
    summary path already truncates to `summarize_doc_max_chars` anyway.
    """
    try:
        import fitz  # type: ignore  # PyMuPDF

        out: list[bytes] = []
        with fitz.open(stream=raw, filetype="pdf") as doc:
            for page_idx in range(min(len(doc), max_pages)):
                page = doc[page_idx]
                # 2x zoom for OCR readability without ballooning memory.
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                out.append(pix.tobytes("png"))
        return out
    except Exception:
        logger.exception("pymupdf: PDF rasterize failed")
        return []


def _pdf_ocr_fallback(raw: bytes) -> str:
    """Last-resort PDF text extraction for image-only PDFs.

    Walks the first 10 pages, rasterizes each via PyMuPDF, runs
    `_ocr_image` (Florence-2 <OCR>) on the raster, concatenates the
    page-level OCR output with double newlines between pages.
    Returns "" if PyMuPDF or Florence-2 are unavailable.
    """
    pages = _pdf_pages_to_images(raw, max_pages=10)
    if not pages:
        return ""
    chunks: list[str] = []
    for png_bytes in pages:
        try:
            page_text = _ocr_image(png_bytes)
        except Exception:
            page_text = None
        if page_text:
            chunks.append(page_text.strip())
    return "\n\n".join(chunks)


def _extract_docx(raw: bytes) -> tuple[str, Optional[str]]:
    import docx  # type: ignore

    doc = docx.Document(BytesIO(raw))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)

    topic: Optional[str] = None
    # Heading 1 wins; otherwise first paragraph.
    for p in doc.paragraphs:
        if p.style and p.style.name and p.style.name.startswith("Heading 1"):
            if p.text.strip():
                topic = p.text.strip()
                break
    if not topic and paragraphs:
        topic = paragraphs[0][:120].strip() or None
    return text, topic


def _extract_xlsx(raw: bytes) -> tuple[str, Optional[str]]:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(BytesIO(raw), read_only=True, data_only=True)
    rows_out: list[str] = []
    sheet_names = wb.sheetnames
    for name in sheet_names[:3]:
        ws = wb[name]
        rows_out.append(f"# Sheet: {name}")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 200:
                break
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                rows_out.append(" | ".join(cells))
    text = "\n".join(rows_out)
    topic = sheet_names[0] if sheet_names else None
    return text, topic


def _extract_pptx(raw: bytes) -> tuple[str, Optional[str]]:
    """Extract slide titles, body text, and speaker notes from a .pptx.

    python-pptx reads OOXML PowerPoint only; a legacy binary .ppt raises,
    which the caller's `except` turns into the harmless filename stub.
    """
    from pptx import Presentation  # type: ignore

    prs = Presentation(BytesIO(raw))
    lines: list[str] = []
    topic: Optional[str] = None
    for idx, slide in enumerate(prs.slides):
        title_txt = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title_txt = slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            title_txt = ""
        lines.append(
            f"# Slide {idx + 1}: {title_txt}" if title_txt
            else f"# Slide {idx + 1}"
        )
        if title_txt and topic is None:
            topic = title_txt[:120]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt and txt != title_txt:
                    lines.append(txt)
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"(notes) {notes}")
        except (AttributeError, ValueError):
            pass
    text = "\n".join(lines)
    if topic is None:
        for ln in lines:
            if not ln.startswith("# Slide"):
                topic = ln[:120]
                break
    return text, topic


def _extract_plaintext(raw: bytes, filename: str) -> tuple[str, Optional[str]]:
    text = raw.decode("utf-8", errors="replace")
    topic: Optional[str] = None
    for line in text.splitlines():
        stripped = line.strip().lstrip("#").strip()
        if 4 <= len(stripped) <= 120:
            topic = stripped
            break
    if not topic:
        topic = filename.rsplit(".", 1)[0]
    return text, topic


def _extract_code_source(filename: str, raw: bytes) -> str:
    """Decode a source-code file to text for the code summarizer.

    UTF-8 with replacement so a stray non-UTF-8 byte never crashes the
    job. `.ipynb` is JSON — pull just the source cells (markdown + code)
    so the summarizer sees the notebook's actual content, not the JSON
    envelope. Everything else is read as plain text and capped so a
    100k-line generated file doesn't blow the regex scans (the head is
    where imports / top-level defs / the module docstring live anyway).
    """
    name = (filename or "").lower()
    if name.endswith(".ipynb"):
        try:
            import json as _json
            nb = _json.loads(raw.decode("utf-8", errors="replace"))
            parts: list[str] = []
            for cell in nb.get("cells", [])[:200]:
                src = cell.get("source") or []
                if isinstance(src, list):
                    parts.append("".join(src))
                elif isinstance(src, str):
                    parts.append(src)
            joined = "\n\n".join(p for p in parts if p.strip())
            if joined.strip():
                return joined[:60000]
        except Exception:
            logger.exception("ipynb parse failed for %s", filename)
        # Fall through to raw text on any parse failure.
    text = raw.decode("utf-8", errors="replace")
    # Cap — the head carries the signal; a giant minified bundle or
    # generated lockfile would otherwise feed 5 MB into the regex scans.
    return text[:120000]


# --- extractive summarization ---------------------------------------------


def _extractive_summary(text: str) -> str:
    """Abstractive document summary via DistilBART CNN.

    Tries BART first (real paraphrase, "The report covers Q4 growth and the
    Phase 11 release"). Falls back to sumy LSA, then to leading sentences,
    so summarization keeps working even if the model can't load.
    """
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return ""

    bart = _bart_summary(text)
    if bart:
        return bart

    n = settings.summarize_doc_sentence_count
    try:
        from sumy.parsers.plaintext import PlaintextParser  # type: ignore
        from sumy.nlp.tokenizers import Tokenizer  # type: ignore
        from sumy.summarizers.lsa import LsaSummarizer  # type: ignore

        parser = PlaintextParser.from_string(text, Tokenizer("english"))
        summarizer = LsaSummarizer()
        sentences = summarizer(parser.document, n)
        joined = " ".join(str(s) for s in sentences).strip()
        if joined:
            return joined
    except Exception:
        logger.exception("sumy LSA failed; using leading sentences")

    parts = re.split(r"(?<=[.!?])\s+", text)
    return " ".join(parts[:n]).strip()


def _llm_doc_summary(
    text: str, filename: str, topic: Optional[str]
) -> Optional[str]:
    """Summarize a document with Qwen2.5-Instruct.

    For docs that fit in one prompt window (~6 000 chars / ~1 500 tokens
    of context for the LLM, leaving room for instructions and reply) we
    summarize directly. Longer docs get chunked, each chunk is
    micro-summarized in one pass, then the chunk summaries are merged in
    a second pass — classic map-reduce. The merged summary captures
    breadth instead of skimming the first page (which was DistilBART's
    failure mode).

    Returns None when the LLM is disabled or unavailable; caller falls
    back to BART → sumy → leading sentences.
    """
    if not settings.rewriter_enabled:
        return None
    text = text.strip()
    if not text:
        return None

    # Qwen2.5-1.5B handles ~3-4 K tokens cleanly; ~6 000 chars is a safe
    # English budget that leaves room for the instruction and reply. Above
    # that we chunk-and-merge.
    CHUNK_BUDGET = 6000
    chunks = _split_for_summary(text, CHUNK_BUDGET)

    if len(chunks) == 1:
        return _llm_summarize_chunk(chunks[0], filename, topic, is_full=True)

    # Map: per-chunk summaries.
    partials: list[str] = []
    for c in chunks[:8]:  # cap fan-out — diminishing returns past 8 chunks
        s = _llm_summarize_chunk(c, filename, topic, is_full=False)
        if s:
            partials.append(s)
    if not partials:
        return None

    # Reduce: condense the partials into one summary.
    merged_input = "\n\n".join(f"- {p}" for p in partials)
    return _llm_merge_summaries(merged_input, filename, topic)


def split_doc_for_embedding(text: str, budget_chars: int = 2000) -> list[str]:
    """Split a document into ~500-token chunks suitable for per-chunk
    CLIP embedding (Sprint I D2). CLIP text encoder caps at 77 tokens
    but truncates gracefully — we keep chunks at 1500-2000 chars so
    each one captures a coherent paragraph or two without losing the
    leading topic sentence to truncation.

    Smaller than the summary chunk budget (6 000 chars) because the
    embedding chunks need to be retrievable as standalone jump-to
    targets — search returns (image_id, chunk_index, snippet) where
    the snippet is the chunk text, so we want enough context for the
    user to recognize but not a whole page-worth.
    """
    return _split_for_summary(text, budget_chars)


def _split_for_summary(text: str, budget_chars: int) -> list[str]:
    """Split on paragraph boundaries when possible, falling back to fixed-
    size windows. Keeps each chunk ≤ budget_chars."""
    if len(text) <= budget_chars:
        return [text]
    paras = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    current = ""
    for p in paras:
        if not p.strip():
            continue
        if len(current) + len(p) + 2 <= budget_chars:
            current = (current + "\n\n" + p) if current else p
        else:
            if current:
                chunks.append(current)
            if len(p) <= budget_chars:
                current = p
            else:
                # Single paragraph longer than the budget — hard-split.
                for i in range(0, len(p), budget_chars):
                    chunks.append(p[i : i + budget_chars])
                current = ""
    if current:
        chunks.append(current)
    return chunks


def _llm_summarize_chunk(
    chunk: str, filename: str, topic: Optional[str], is_full: bool
) -> Optional[str]:
    """One LLM call. is_full=True asks for a final-quality summary; False
    asks for a terse partial that will be merged later."""
    try:
        import torch

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        if is_full:
            instructions = (
                "Write a SHORT, search-friendly description of what this "
                "document IS — the kind of thing someone would type into "
                "a search bar months later to find it by memory. "
                "1-2 natural English sentences, under 35 words total. "
                "Describe the document's PURPOSE and SUBJECT (what it's "
                "about, what kind of doc it is), NOT its body content. "
                "Pack in the most distinctive concrete nouns — project "
                "names, parties, topics, dates — but DO NOT quote or "
                "list the document's questions, prompts, or paragraphs "
                "verbatim. NEVER copy a sentence from the input — "
                "always paraphrase. Do NOT begin with 'This document', "
                "'The document', or 'In summary'. Output only the "
                "description.\n\n"
                "Examples of GOOD descriptions:\n"
                "- 'Class assignment asking students to inventory what "
                "they built or trained during the course.'\n"
                "- 'Group exercise comparing four AI agent frameworks "
                "(LangChain, AutoGen, CrewAI, Semantic Kernel) for a "
                "research write-up.'\n"
                "- 'Q4 2025 financial report covering revenue growth, "
                "operating margin, and Phase 11 product release plans.'\n"
                "Examples of BAD descriptions (verbatim body, do NOT "
                "produce these):\n"
                "- 'Part 1 — Inventory Think about or list everything "
                "you can remember…'\n"
                "- 'Phase 1: Framework Research (10 minutes) Your "
                "group has been assigned…'"
            )
        else:
            instructions = (
                "Briefly state what this section of a document is ABOUT "
                "in ONE specific English sentence (under 25 words). "
                "Describe its subject/purpose, do NOT quote body text "
                "verbatim. Output only the sentence."
            )

        ctx_lines = []
        if filename:
            ctx_lines.append(f"Filename: {filename}")
        if topic:
            ctx_lines.append(f"Title: {topic}")
        ctx_lines.append("Content:\n" + chunk)

        messages = [
            {
                "role": "system",
                "content": (
                    "You write concise, factual document summaries. "
                    "You never invent details that are not in the input."
                ),
            },
            {
                "role": "user",
                "content": instructions + "\n\n" + "\n".join(ctx_lines),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3500
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]

        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=240 if is_full else 120,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        # Strip surrounding quotes / leading dashes.
        reply = reply.lstrip("-•* ").strip().strip('"').strip("'").strip()
        if not reply:
            return None
        # Keep the reply within sane bounds; the LLM occasionally rambles.
        if len(reply) > 600:
            reply = reply[:600].rsplit(".", 1)[0] + "."
        return reply
    except Exception:
        logger.exception("qwen: doc summary chunk failed")
        return None


def _llm_merge_summaries(
    partials_block: str, filename: str, topic: Optional[str]
) -> Optional[str]:
    try:
        import torch

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        ctx_lines = []
        if filename:
            ctx_lines.append(f"Filename: {filename}")
        if topic:
            ctx_lines.append(f"Title: {topic}")
        ctx_lines.append(
            "Partial summaries (each line summarizes a different "
            "section of the same document):"
        )
        ctx_lines.append(partials_block)

        instructions = (
            "Combine the partial summaries below into ONE SHORT "
            "search-friendly description of what the document IS — "
            "the kind of thing someone would type into a search bar "
            "months later to find it. 1-2 natural English sentences, "
            "under 35 words total. Describe purpose and subject, drop "
            "redundancy and any verbatim body content. Keep the most "
            "distinctive specific nouns (names, parties, topics, dates). "
            "Do NOT begin with 'This document', 'The document', or 'In "
            "summary'. Output only the combined description."
        )

        messages = [
            {
                "role": "system",
                "content": (
                    "You write concise, factual document summaries. "
                    "You never invent details."
                ),
            },
            {
                "role": "user",
                "content": instructions + "\n\n" + "\n".join(ctx_lines),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3500
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]

        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=240,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        reply = reply.lstrip("-•* ").strip().strip('"').strip("'").strip()
        if not reply:
            return None
        if len(reply) > 600:
            reply = reply[:600].rsplit(".", 1)[0] + "."
        return reply
    except Exception:
        logger.exception("qwen: doc merge failed")
        return None


def _llm_doc_topic(text: str, filename: str) -> Optional[str]:
    """Ask the LLM for a 3–8 word topic line when extraction failed."""
    try:
        import torch

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        head = text[:2000]
        messages = [
            {
                "role": "system",
                "content": (
                    "You write concise, factual document titles. "
                    "You never invent details."
                ),
            },
            {
                "role": "user",
                "content": (
                    "Give a 3–8 word title that captures what this "
                    "document is about. Use noun phrases (no verbs, no "
                    "punctuation at the end). Output only the title.\n\n"
                    f"Filename: {filename}\n\nContent:\n{head}"
                ),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3000
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]
        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=24,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        reply = reply.split("\n")[0].strip().strip('"').strip("'").strip()
        if 3 <= len(reply) <= 80:
            return reply
        return None
    except Exception:
        logger.exception("qwen: doc topic failed")
        return None


def _llm_keypoints(text: str, filename: str) -> list[str]:
    """Ask the LLM for 3 short key points when the heuristic returned
    nothing (i.e. the doc has no markdown list syntax)."""
    try:
        import torch

        from backend.vision.runtime import get_summary_rewriter

        model, tokenizer, device = get_summary_rewriter()

        head = text[:6000]
        messages = [
            {
                "role": "system",
                "content": (
                    "You extract concrete key points from documents. "
                    "You never invent details."
                ),
            },
            {
                "role": "user",
                "content": (
                    "List up to 3 KEY POINTS from this document, one per "
                    "line. Each point: a short clause (under 12 words) "
                    "covering a specific fact, decision, or section. "
                    "Output ONLY the lines, no numbering, no bullets, "
                    "no preamble.\n\n"
                    f"Filename: {filename}\n\nContent:\n{head}"
                ),
            },
        ]
        prompt = tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        inputs = tokenizer(
            prompt, return_tensors="pt", truncation=True, max_length=3500
        ).to(device)
        prompt_len = inputs.input_ids.shape[1]
        with torch.no_grad():
            out_ids = model.generate(
                **inputs,
                max_new_tokens=120,
                do_sample=False,
                num_beams=1,
                pad_token_id=tokenizer.pad_token_id or tokenizer.eos_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        new_ids = out_ids[0][prompt_len:]
        reply = tokenizer.decode(new_ids, skip_special_tokens=True).strip()
        points = []
        for line in reply.splitlines():
            cleaned = line.strip().lstrip("-•*").strip()
            cleaned = re.sub(r"^\d+[\.\)]\s*", "", cleaned)
            if 4 <= len(cleaned) <= 160:
                points.append(cleaned)
            if len(points) >= 3:
                break
        return points
    except Exception:
        logger.exception("qwen: keypoints failed")
        return []


def _bart_summary(text: str) -> Optional[str]:
    """Run DistilBART abstractive summarization. Returns None on any failure
    (model unavailable, OOM, tokenizer error). Caller falls back to sumy.

    BART has a 1024-token input cap; we truncate to ~3500 chars (≈ 800
    tokens for English) to leave room for the model's encoder context.
    """
    try:
        import torch

        from backend.vision.runtime import get_doc_summarizer

        model, tokenizer, device = get_doc_summarizer()
        truncated = text[:3500]
        inputs = tokenizer(
            truncated,
            return_tensors="pt",
            truncation=True,
            max_length=1024,
        ).to(device)

        with torch.no_grad():
            ids = model.generate(
                **inputs,
                max_length=140,
                min_length=40,
                num_beams=4,
                length_penalty=2.0,
                no_repeat_ngram_size=3,
                early_stopping=True,
            )
        out = tokenizer.decode(ids[0], skip_special_tokens=True).strip()
        return out or None
    except Exception:
        logger.exception("bart: summary failed")
        return None


def _extract_keypoints(text: str) -> list[str]:
    """Cheap keypoint heuristic: list-formatted lines + frequent capitalized
    nouns. Better than nothing without spinning up an LLM."""
    points: list[str] = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        m = re.match(r"^[\-\*•–—]\s+(.{4,160})$", line)
        if m:
            points.append(m.group(1).strip())
        elif re.match(r"^\d+[\.\)]\s+(.{4,160})$", line):
            points.append(re.sub(r"^\d+[\.\)]\s+", "", line).strip())
        if len(points) >= 5:
            break
    return points
